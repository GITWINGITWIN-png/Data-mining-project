"""สร้างสไลด์นำเสนอ (.pptx) ที่เปิด/แก้ต่อได้ใน Google Slides

    python make_slides.py            # เขียน ElderCare_Insight_Slides.pptx ข้าง ๆ ไฟล์นี้

ทำไมต้อง generate ไม่ใช่วาดมือ: ตัวเลขบนสไลด์ทุกตัวมาจากคลังข้อมูลชุดเดียวกับที่รายงาน
และแดชบอร์ดใช้ ถ้าคลังถูก build ใหม่แล้วตัวเลขขยับ แก้ที่ค่าคงที่ด้านล่างที่เดียวแล้ว
รันซ้ำ สไลด์กับรายงานจะไม่หลุดจากกันเงียบ ๆ เหมือนตอนที่แก้ทีละกล่องข้อความ

ต้องมี python-pptx:  uv run --with python-pptx python make_slides.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
DASH = HERE.parent / "04_Dashboard"
SHOTS = DASH / "screenshots"
OUT = HERE / "ElderCare_Insight_Slides.pptx"

# ---------------------------------------------------------------------------
# จานสี — ตรงกับ 04_Dashboard/charts.py ทุกค่า เพื่อให้สไลด์กับกราฟที่วางลงไป
# ดูเป็นระบบเดียวกัน ไม่ใช่ภาพที่ถูกแปะมาจากที่อื่น
# ---------------------------------------------------------------------------
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
PANEL = RGBColor(0xF4, 0xF4, 0xF0)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
GRID = RGBColor(0xE1, 0xE0, 0xD9)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
ACCENT_TINT = RGBColor(0xEA, 0xF1, 0xFB)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
GREEN = RGBColor(0x1B, 0xAF, 0x7A)
WARN = RGBColor(0xC9, 0x42, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEAD_FONT = "Prompt"      # ฟอนต์ไทยแบบไม่มีหัว ใช้กับหัวเรื่องและตัวเลขใหญ่
BODY_FONT = "Sarabun"     # ฟอนต์ไทยมาตรฐานสำหรับเนื้อความ

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M = Inches(0.72)                       # ขอบซ้าย/ขวา
CONTENT_W = SLIDE_W - 2 * M
BODY_TOP = Inches(1.62)                # ใต้เส้นคั่นหัวเรื่อง
BODY_BOTTOM = Inches(6.86)             # เหนือ footer
BODY_H = BODY_BOTTOM - BODY_TOP

DECK_NAME = "ElderCare Insight · Mini DW & Dashboard"

_page = [0]


# ---------------------------------------------------------------------------
# helper ระดับล่าง
# ---------------------------------------------------------------------------
def _set_font(run, *, size, bold=False, color=INK, font=BODY_FONT, italic=False):
    """ตั้งฟอนต์ให้ครบทั้งสามช่อง latin / ea / cs

    ตัวอักษรไทยถูกจัดเป็น "complex script" ถ้าตั้งแต่ latin อย่างเดียว PowerPoint
    จะไปหยิบฟอนต์ cs ของธีมมาใช้แทน แล้วสไลด์จะดูไม่เหมือนที่ตั้งใจ
    """
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:cs", "a:ea"):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
        el = parse_xml(
            '<a:%s xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'typeface="%s"/>' % (tag.split(":")[1], font)
        )
        rPr.append(el)


def _rich(p, text, *, size, color=INK, font=BODY_FONT, bold=False):
    """เขียนข้อความลงย่อหน้า โดยแปลง **...** เป็นตัวหนา

    ตัวเลขคือสิ่งที่คนอ่านสไลด์ต้องจับให้ได้ก่อน จึงคุ้มที่จะมี markup เล็ก ๆ
    แทนที่จะต้องแตก run เองทุกที่
    """
    for i, chunk in enumerate(text.split("**")):
        if not chunk:
            continue
        run = p.add_run()
        run.text = chunk
        _set_font(run, size=size, bold=bold or bool(i % 2), color=color, font=font)


def textbox(slide, left, top, width, height, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size, color=INK, font=BODY_FONT, bold=False, first=False,
         space_before=0, space_after=6, line_spacing=1.18, align=PP_ALIGN.LEFT):
    """เขียนข้อความหนึ่งบล็อก โดยแตก \n เป็นย่อหน้าใหม่

    run ของ pptx ใส่ newline ไม่ได้ (มันจะกลายเป็นอักขระควบคุมที่ PowerPoint ปฏิเสธ)
    บรรทัดว่างจึงถูกใช้เป็นตัวคั่นย่อหน้าและแปลงเป็นระยะห่างแทน
    """
    p = None
    written = 0
    extra_gap = 0
    for part in text.split("\n"):
        if not part.strip():
            extra_gap = 5
            continue
        p = tf.paragraphs[0] if (first and written == 0) else tf.add_paragraph()
        written += 1
        p.space_before = Pt(space_before + extra_gap)
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        p.alignment = align
        _rich(p, part, size=size, color=color, font=font, bold=bold)
        extra_gap = 0
    return p


def rect(slide, left, top, width, height, *, fill=None, line=None, line_w=0.75,
         shape=MSO_SHAPE.RECTANGLE, radius=None):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = radius
    sh.text_frame.word_wrap = True
    return sh


def hline(slide, left, top, width, *, color=GRID, w=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    return ln


def place_image(slide, path, left, top, width, height, *, border=True):
    """วางภาพให้พอดีกรอบโดยรักษาสัดส่วน แล้วจัดกึ่งกลางในกรอบนั้น"""
    pic = slide.shapes.add_picture(str(path), left, top)
    box_ar = width / height
    ar = pic.width / pic.height
    if ar >= box_ar:
        w, h = width, int(width / ar)
    else:
        h, w = height, int(height * ar)
    pic.width, pic.height = int(w), int(h)
    pic.left = int(left + (width - pic.width) / 2)
    pic.top = int(top + (height - pic.height) / 2)
    if border:
        pic.line.color.rgb = GRID
        pic.line.width = Pt(0.75)
    return pic


# ---- ตาราง ---------------------------------------------------------------
_NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # "No Style, No Grid"


def _cell_border(cell, edge, color, width_pt):
    """ใส่เส้นขอบให้เซลล์

    ลำดับลูกใน a:tcPr ตาม schema คือ lnL, lnR, lnT, lnB แล้วค่อยถึง fill
    จึงต้อง insert ไว้หน้าสุดเสมอ ไม่งั้น PowerPoint จะปฏิเสธไฟล์
    """
    tcPr = cell._tc.get_or_add_tcPr()
    tag = "a:ln" + edge
    for el in tcPr.findall(qn(tag)):
        tcPr.remove(el)
    el = parse_xml(
        '<a:ln%s xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'w="%d" cap="flat" cmpd="sng" algn="ctr">'
        '<a:solidFill><a:srgbClr val="%02X%02X%02X"/></a:solidFill></a:ln%s>'
        % (edge, int(width_pt * 12700), color[0], color[1], color[2], edge)
    )
    tcPr.insert(0, el)


def table(slide, rows_data, left, top, width, col_widths, *,
          header_size=10.5, body_size=10.0, row_h=Inches(0.34),
          header_fill=ACCENT_TINT, aligns=None):
    """ตารางแบนราบ: หัวตารางเป็นสีฟ้าจาง แถวสลับสีอ่อน เส้นคั่นบาง ๆ อย่างเดียว"""
    n_rows, n_cols = len(rows_data), len(rows_data[0])
    gf = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows)
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False
    tblPr = tbl._tbl.tblPr
    for el in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(el)
    sid = parse_xml(
        '<a:tableStyleId xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        "%s</a:tableStyleId>" % _NO_STYLE
    )
    tblPr.append(sid)

    total = sum(col_widths)
    for i, frac in enumerate(col_widths):
        tbl.columns[i].width = Emu(int(width * frac / total))
    for r in range(n_rows):
        tbl.rows[r].height = row_h

    aligns = aligns or [PP_ALIGN.LEFT] * n_cols
    for r, row in enumerate(rows_data):
        head = r == 0
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            _cell_border(cell, "B", GRID if not head else RGBColor(0xC3, 0xC2, 0xB7), 0.75)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                header_fill if head else (WHITE if r % 2 else PANEL)
            )
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = aligns[c]
            p.line_spacing = 1.06
            _rich(
                p, str(value),
                size=header_size if head else body_size,
                color=INK if head else INK2,
                font=HEAD_FONT if head else BODY_FONT,
                bold=head,
            )
    return tbl


# ---------------------------------------------------------------------------
# โครงสไลด์
# ---------------------------------------------------------------------------
def new_slide(prs, *, bg=SURFACE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def footer(slide):
    _page[0] += 1
    tf = textbox(slide, M, Inches(6.98), CONTENT_W, Inches(0.26))
    p = tf.paragraphs[0]
    _rich(p, DECK_NAME, size=8.5, color=MUTED)
    tf2 = textbox(slide, SLIDE_W - M - Inches(1.2), Inches(6.98), Inches(1.2), Inches(0.26))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _rich(p2, str(_page[0]), size=8.5, color=MUTED, font=HEAD_FONT)


def head(slide, kicker, title, *, sub=None):
    """หัวเรื่องมาตรฐาน: kicker เล็กสีฟ้า · ชื่อสไลด์ · เส้นคั่น"""
    tf = textbox(slide, M, Inches(0.44), CONTENT_W, Inches(0.24))
    _rich(tf.paragraphs[0], kicker, size=10.5, color=ACCENT, font=HEAD_FONT, bold=True)
    tf2 = textbox(slide, M, Inches(0.70), CONTENT_W, Inches(0.55))
    p = tf2.paragraphs[0]
    p.line_spacing = 1.0
    _rich(p, title, size=25, color=INK, font=HEAD_FONT, bold=True)
    y = Inches(1.34)
    if sub:
        tf3 = textbox(slide, M, Inches(1.24), CONTENT_W, Inches(0.3))
        _rich(tf3.paragraphs[0], sub, size=12.5, color=INK2)
        y = Inches(1.62)
        hline(slide, M, y - Inches(0.14), CONTENT_W)
        return y
    hline(slide, M, y, CONTENT_W)
    return BODY_TOP


def bullets(tf, items, *, size=13.5, color=INK2, gap=9, bullet_color=ACCENT, first=False):
    """รายการหัวข้อย่อย — จุดนำเป็นสีเน้น ตัวข้อความเป็นสีหมึกรอง"""
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if (first and i == 0) else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.2
        b = p.add_run()
        b.text = "•  "
        _set_font(b, size=size, bold=True, color=bullet_color, font=BODY_FONT)
        _rich(p, item, size=size, color=color)


def callout(slide, left, top, width, height, title, body, *,
            accent=ACCENT, fill=ACCENT_TINT, title_size=12, body_size=11.5):
    """กล่องเน้น — แถบสีด้านซ้าย พื้นจาง ใช้กับ 'กับดักที่ปิดไปแล้ว' เป็นหลัก"""
    rect(slide, left, top, width, height, fill=fill)
    rect(slide, left, top, Inches(0.055), height, fill=accent)
    tf = textbox(slide, left + Inches(0.26), top + Inches(0.13),
                 width - Inches(0.44), height - Inches(0.26))
    if title:
        para(tf, title, size=title_size, color=accent, font=HEAD_FONT, bold=True,
             first=True, space_after=4)
        para(tf, body, size=body_size, color=INK2, space_after=0)
    else:
        para(tf, body, size=body_size, color=INK2, first=True, space_after=0)


def stat_bar(slide, left, top, width, height, value, label, note, *, color=ACCENT):
    """การ์ดสถิติแนวนอน — ตัวเลขซ้าย คำอธิบายขวา ใช้ตอนที่ความสูงมีจำกัด"""
    rect(slide, left, top, width, height, fill=WHITE, line=GRID)
    rect(slide, left, top, Inches(0.05), height, fill=color)
    tf = textbox(slide, left + Inches(0.24), top, Inches(1.6), height,
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, value, size=19, color=color, font=HEAD_FONT, bold=True, first=True,
         space_after=0, line_spacing=1.0)
    tf2 = textbox(slide, left + Inches(1.98), top, width - Inches(2.2), height,
                  anchor=MSO_ANCHOR.MIDDLE)
    para(tf2, label, size=11.5, color=INK, bold=True, first=True, space_after=1)
    para(tf2, note, size=9.5, color=MUTED, space_after=0)


def stat_card(slide, left, top, width, height, value, label, note=None, *, color=ACCENT):
    rect(slide, left, top, width, height, fill=WHITE, line=GRID)
    rect(slide, left, top, width, Inches(0.05), fill=color)
    tf = textbox(slide, left + Inches(0.2), top + Inches(0.24),
                 width - Inches(0.4), height - Inches(0.4))
    para(tf, value, size=26, color=color, font=HEAD_FONT, bold=True, first=True,
         space_after=2, line_spacing=1.0)
    para(tf, label, size=11.5, color=INK, bold=True, space_after=2)
    if note:
        para(tf, note, size=9.5, color=MUTED, space_after=0)


# ---------------------------------------------------------------------------
# ไดอะแกรม — วาดด้วยรูปทรงของ pptx เอง ไม่ใช่ภาพ เพื่อให้แก้ต่อใน Google Slides ได้
# ---------------------------------------------------------------------------
def _box(slide, left, top, width, height, title, lines, *, accent, title_size=10.5,
         line_size=7.6):
    rect(slide, left, top, width, height, fill=WHITE, line=GRID)
    rect(slide, left, top, width, Inches(0.045), fill=accent)
    tf = textbox(slide, left + Inches(0.12), top + Inches(0.15),
                 width - Inches(0.24), height - Inches(0.24))
    para(tf, title, size=title_size, color=accent, font=HEAD_FONT, bold=True,
         first=True, space_after=3, line_spacing=1.0)
    for ln in lines:
        para(tf, ln, size=line_size, color=INK2, space_after=1, line_spacing=1.08)


def star_schema(slide, top):
    """Double star: Fact สองตัวเรียงกลาง · Dimension หกตัวขนาบซ้ายขวา"""
    lw, cw, rw = Inches(2.85), Inches(4.35), Inches(2.85)
    lx = M
    cx = M + lw + Inches(0.95)
    rx = cx + cw + Inches(0.95)
    dim_h = Inches(1.30)
    ys = [top + Inches(0.02), top + Inches(1.62), top + Inches(3.22)]

    left_dims = [
        ("Dim_Date", ["PK date_key", "full_date · year · quarter · month", "year_month · us_fiscal_year", "covid_period"]),
        ("Dim_Geography", ["PK geography_key", "zip_code · city · county_parish", "state_code · state_name", "census_region · urban_rural"]),
        ("Dim_Ownership", ["PK ownership_key", "ownership_type (13 ประเภท)", "ownership_group (3 กลุ่ม)", "is_for_profit"]),
    ]
    right_dims = [
        ("Dim_Facility  (SCD2)", ["PK facility_key", "ccn · provider_name · provider_type", "special_focus_status · abuse_icon", "effective_date · expiry_date · is_current"]),
        ("Dim_Chain", ["PK chain_key", "chain_id · chain_name", "num_facilities_in_chain", "chain_size_band · is_independent"]),
        ("Dim_Penalty_Type", ["PK penalty_type_key", "penalty_type", "is_monetary", "ใช้กับ Fact 2 เท่านั้น"]),
    ]

    boxes_l, boxes_r = [], []
    for (name, lines), y in zip(left_dims, ys):
        _box(slide, lx, y, lw, dim_h, name, lines, accent=ACCENT)
        boxes_l.append((lx, y, lw, dim_h))
    for (name, lines), y in zip(right_dims, ys):
        _box(slide, rx, y, rw, dim_h, name, lines, accent=ACCENT)
        boxes_r.append((rx, y, rw, dim_h))

    f1y, f2y = top + Inches(0.18), top + Inches(2.58)
    fh = Inches(1.92)
    _box(slide, cx, f1y, cw, fh, "Fact_Facility_Monthly   ·   periodic snapshot",
         ["Grain — สถานพยาบาลหนึ่งแห่ง (หนึ่ง CCN) ณ วันประมวลผลหนึ่งงวด",
          "PK (snapshot_date_key, ccn)     DD ccn",
          "FK snapshot_date_key · facility_key · geography_key · ownership_key · chain_key",
          "— measures — certified_beds · avg_residents_per_day · resident_days",
          "reported_total_nurse_hprd · reported_rn_hprd · total_nursing_turnover_pct",
          "overall_rating · staffing_rating · health_inspection_rating",
          "cycle1_total_deficiencies · is_suspect"],
         accent=ORANGE, title_size=11, line_size=7.8)
    _box(slide, cx, f2y, cw, fh, "Fact_Penalty_Event   ·   transaction",
         ["Grain — การลงโทษหนึ่งครั้งต่อสถานพยาบาลหนึ่งแห่งในวันหนึ่ง",
          "PK penalty_event_key     DD ccn · fine_id",
          "FK penalty_date_key · facility_key · geography_key · ownership_key",
          "chain_key · penalty_type_key",
          "— measures — fine_amount_usd · payment_denial_days",
          "penalty_count (= 1) · fine_id_source"],
         accent=ORANGE, title_size=11, line_size=7.8)

    def link(x1, y1, x2, y2, dashed=False):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
        ln.line.color.rgb = RGBColor(0xC3, 0xC2, 0xB7)
        ln.line.width = Pt(0.75)
        if dashed:
            ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        return ln

    # dimension ทั้งห้าตัวแรกเป็น conformed — มีเส้นไปหา Fact ทั้งสองตัว
    for i, (x, y, w, h) in enumerate(boxes_l):
        link(x + w, y + h / 2, cx, f1y + fh * (0.3 + 0.2 * i))
        link(x + w, y + h / 2, cx, f2y + fh * (0.3 + 0.2 * i))
    for i, (x, y, w, h) in enumerate(boxes_r[:2]):
        link(x, y + h / 2, cx + cw, f1y + fh * (0.32 + 0.24 * i))
        link(x, y + h / 2, cx + cw, f2y + fh * (0.32 + 0.24 * i))
    x, y, w, h = boxes_r[2]
    link(x, y + h / 2, cx + cw, f2y + fh * 0.82)


def etl_diagram(slide, top):
    stages = [
        ("Extract", "fetch_snapshots.py", ACCENT,
         ["ไล่รายการ 88 งวดผ่าน API", "ดาวน์โหลดแบบ incremental", "แตกเฉพาะ 6 CSV ที่ใช้จริง"]),
        ("Clean", "clean.py · columns.py", ORANGE,
         ["กฎ Q1, Q3, Q5, Q6, Q7", "แม็ปคอลัมน์รายยุค 78 → 99", "รหัสเชิงอรรถออกจากตัวหาร"]),
        ("Transform", "dimensions.py · facts.py", ORANGE,
         ["6 Dimension รวม SCD2", "2 Fact + กฎ Q2, Q4, Q8", "คำนวณ M1–M8"]),
        ("Integrate", "sources.py · population.py", ORANGE,
         ["กระทบยอดงวดที่เผยแพร่ซ้ำ", "เชื่อมด้วย ccn ทุกจุด", "ต่อประชากร 65+ ที่ (รัฐ, ปี)"]),
        ("Load", "load.py · build_warehouse.py", GREEN,
         ["ตรวจ PK/FK/SCD2 ก่อนเขียน", "สร้างตารางใหม่จาก schema.sql", "ใส่ข้อมูลผ่าน constraint"]),
    ]
    n = len(stages)
    gap = Inches(0.18)
    w = int((CONTENT_W - gap * (n - 1)) / n)
    h = Inches(0.92)
    for i, (name, module, color, lines) in enumerate(stages):
        x = M + i * (w + gap)
        rect(slide, x, top, w, h, fill=color)
        tf = textbox(slide, x + Inches(0.14), top + Inches(0.14), w - Inches(0.28),
                     h - Inches(0.28))
        para(tf, name, size=15, color=WHITE, font=HEAD_FONT, bold=True, first=True,
             space_after=1, line_spacing=1.0)
        para(tf, module, size=8.5, color=WHITE, space_after=0)
        if i < n - 1:
            ar = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE, x + w + Inches(0.025),
                top + h / 2 - Inches(0.075), Inches(0.13), Inches(0.15))
            ar.rotation = 90
            ar.shadow.inherit = False
            ar.fill.solid()
            ar.fill.fore_color.rgb = MUTED
            ar.line.fill.background()
        card_top = top + h + Inches(0.16)
        rect(slide, x, card_top, w, Inches(1.28), fill=WHITE, line=GRID)
        tf2 = textbox(slide, x + Inches(0.14), card_top + Inches(0.15),
                      w - Inches(0.28), Inches(1.0))
        for j, ln in enumerate(lines):
            para(tf2, "· " + ln, size=9, color=INK2, first=(j == 0), space_after=4,
                 line_spacing=1.1)


# ---------------------------------------------------------------------------
# สไลด์
# ---------------------------------------------------------------------------
def slide_title(prs):
    slide = new_slide(prs, bg=INK)
    rect(slide, 0, 0, Inches(0.09), SLIDE_H, fill=ACCENT)
    tf = textbox(slide, Inches(1.15), Inches(1.55), Inches(10.5), Inches(0.3))
    _rich(tf.paragraphs[0], "GROUP ASSIGNMENT #1  ·  DATA MINING / DATA WAREHOUSE",
          size=11.5, color=ACCENT, font=HEAD_FONT, bold=True)

    tf = textbox(slide, Inches(1.15), Inches(2.02), Inches(11), Inches(1.1))
    para(tf, "ElderCare Insight", size=54, color=WHITE, font=HEAD_FONT, bold=True,
         first=True, space_after=0, line_spacing=1.0)

    tf = textbox(slide, Inches(1.15), Inches(3.12), Inches(10.4), Inches(0.9))
    para(tf, "Mini Data Warehouse & Analytics Dashboard",
         size=21, color=RGBColor(0xE1, 0xE0, 0xD9), font=HEAD_FONT, first=True,
         space_after=6, line_spacing=1.0)
    para(tf, "สำหรับธุรกิจสถานพยาบาลผู้สูงอายุ (skilled nursing facility) ในสหรัฐอเมริกา",
         size=14, color=MUTED, space_after=0)

    hline(slide, Inches(1.15), Inches(4.30), Inches(5.4), color=RGBColor(0x3A, 0x39, 0x36))

    tf = textbox(slide, Inches(1.15), Inches(4.52), Inches(5.6), Inches(1.5))
    para(tf, "**กลุ่มที่ 01**", size=13, color=WHITE, first=True, space_after=5)
    for label in ("สมาชิก 1 __________________  ·  หน้าที่ __________",
                  "สมาชิก 2 __________________  ·  หน้าที่ __________",
                  "สมาชิก 3 __________________  ·  หน้าที่ __________"):
        para(tf, label, size=12, color=RGBColor(0xB5, 0xB3, 0xAD), space_after=4)

    tf = textbox(slide, Inches(7.2), Inches(4.52), Inches(5.0), Inches(1.6))
    for label in ("ข้อมูล CMS Provider Data Catalog · 32 งวดรายไตรมาส 2562–2569",
                  "DuckDB star schema · 2 Fact · 6 Dimension · 483,183 แถว",
                  "ETL เขียนด้วย Python + pandas · แดชบอร์ด Streamlit",
                  "22 สิงหาคม 2569"):
        para(tf, label, size=11.5, color=RGBColor(0xB5, 0xB3, 0xAD),
             first=(label.startswith("ข้อมูล")), space_after=6)


def slide_agenda(prs):
    slide = new_slide(prs)
    y = head(slide, "สารบัญ", "โครงการนำเสนอ และผู้พูดแต่ละช่วง",
             sub="แบ่งตามสัดส่วน 2 / 5 / 5 นาทีที่โจทย์กำหนด — สมาชิกทุกคนพูดและอธิบายผลงานของตนเอง")
    parts = [
        ("01", "Business Problem & Requirements", "2 นาที", ACCENT,
         ["ปัญหาทางธุรกิจและบริบท", "Stakeholders 5 กลุ่ม",
          "Business Questions 8 ข้อ", "Measures 10 ตัว และเรื่องการบวกได้"]),
        ("02", "Data Warehouse Design & ETL", "5 นาที", ORANGE,
         ["แหล่งข้อมูล 5 แหล่ง และความซับซ้อน", "ปัญหาคุณภาพข้อมูล Q1–Q8",
          "Star schema · Fact · Grain · SCD2", "ETL 5 ขั้น และการตรวจก่อน/หลัง Load"]),
        ("03", "Dashboard & Business Insights", "5 นาที", GREEN,
         ["สาธิตแดชบอร์ดและ Interactive Filter", "คำตอบของ BQ1–BQ8 พร้อมตัวเลข",
          "ข้อเสนอแนะทางธุรกิจ 8 ข้อ", "การใช้ Generative AI และการตรวจสอบ"]),
    ]
    w = int((CONTENT_W - Inches(0.4)) / 3)
    for i, (num, title, mins, color, items) in enumerate(parts):
        x = M + i * (w + Inches(0.2))
        rect(slide, x, y, w, Inches(4.5), fill=WHITE, line=GRID)
        rect(slide, x, y, w, Inches(0.06), fill=color)
        tf = textbox(slide, x + Inches(0.28), y + Inches(0.32), w - Inches(0.56), Inches(1.3))
        para(tf, num, size=34, color=color, font=HEAD_FONT, bold=True, first=True,
             space_after=2, line_spacing=1.0)
        para(tf, title, size=14, color=INK, font=HEAD_FONT, bold=True, space_after=3,
             line_spacing=1.1)
        para(tf, mins + "  ·  ผู้พูด ____________", size=11, color=MUTED, space_after=0)
        tf2 = textbox(slide, x + Inches(0.28), y + Inches(1.92), w - Inches(0.56), Inches(2.4))
        bullets(tf2, items, size=11.5, gap=8, bullet_color=color, first=True)


def slide_section(prs, num, title, subtitle, color=ACCENT):
    slide = new_slide(prs, bg=INK)
    rect(slide, 0, 0, Inches(0.09), SLIDE_H, fill=color)
    tf = textbox(slide, Inches(1.15), Inches(2.55), Inches(10.5), Inches(1.0))
    para(tf, num, size=15, color=color, font=HEAD_FONT, bold=True, first=True,
         space_after=8, line_spacing=1.0)
    para(tf, title, size=40, color=WHITE, font=HEAD_FONT, bold=True, space_after=10,
         line_spacing=1.0)
    para(tf, subtitle, size=14, color=MUTED, space_after=0)


def slide_problem(prs):
    slide = new_slide(prs)
    y = head(slide, "01 · BUSINESS PROBLEM", "ข้อมูลมีอยู่แล้ว — ปัญหาคือรูปร่างของมัน")
    tf = textbox(slide, M, y, Inches(7.05), Inches(4.4))
    para(tf, "ผู้ประกอบการสถานพยาบาลผู้สูงอายุต้องตัดสินใจสองเรื่องที่ใช้เงินมากที่สุด — "
             "**จะขยายไปรัฐไหน** และ **ควรลงเงินกับกำลังคนเท่าไร**",
         size=14, color=INK, first=True, space_after=12)
    bullets(tf, [
        "CMS เผยแพร่ข้อมูลสถานพยาบาลที่ได้รับการรับรองทุกแห่งในสหรัฐฯ อยู่แล้ว "
        "**ข้อมูลจึงมีอยู่จริง** สิ่งที่ไม่มีคือรูปร่างที่เอาไปถามคำถามได้",
        "CMS เผยแพร่เฉพาะ **สถานะปัจจุบัน** ทุกงวดเขียนทับงวดก่อนหน้า "
        "ประวัติมีอยู่ในรูปแฟ้ม ZIP ที่ถูกเก็บไว้ **88 แฟ้ม** ซึ่งไม่มีใครต่อเป็นชุดเดียว",
        "คำถามง่าย ๆ อย่าง “อัตราการเข้าพักฟื้นหรือยัง และฟื้นที่ไหน” "
        "จึงไม่มีตารางให้ query เลย",
        "คอลัมน์เดียวที่ดูเหมือนประวัติ — Total Amount of Fines in Dollars — "
        "เป็น **หน้าต่างหมุน 3 ปี** มันลดลงเองได้เมื่อค่าปรับเก่าหลุดออกจากหน้าต่าง "
        "อ่านเป็นยอดสะสมจะได้คำตอบที่มั่นใจและผิด",
    ], size=13, gap=11)

    callout(slide, Inches(8.2), y, Inches(4.41), Inches(1.85),
            "โครงงานนี้จึงสร้างสินทรัพย์ที่ขาดหายไป",
            "คลังข้อมูลที่ **เก็บทุกงวดไว้** เพื่อให้ถามคำถามเรื่องการเปลี่ยนแปลง"
            "ตามเวลาได้ตั้งแต่แรก แทนที่จะเดาจากภาพนิ่งภาพเดียว")
    for i, (val, label, note, col) in enumerate([
        ("32", "งวดรายไตรมาส", "2562-01 ถึง 2569-07", ACCENT),
        ("483,183", "แถวใน Fact 1", "16,155 สถานพยาบาล", ORANGE),
        ("79,803", "เหตุการณ์ลงโทษ", "ค่าปรับรวม $1,616,747,274", GREEN),
    ]):
        stat_bar(slide, Inches(8.2), y + Inches(2.10) + i * Inches(0.98),
                 Inches(4.41), Inches(0.86), val, label, note, color=col)


def slide_stakeholders(prs):
    slide = new_slide(prs)
    y = head(slide, "01 · STAKEHOLDERS", "ใครใช้ข้อมูลนี้ และเขาเป็นเจ้าของการตัดสินใจอะไร")
    rows = [
        ["ผู้ใช้", "การตัดสินใจที่เขาเป็นเจ้าของ", "สิ่งที่ต้องการจากข้อมูล", "ตอบด้วย"],
        ["CEO / หัวหน้าฝ่ายขยายกิจการ", "จะเข้าตลาดใด หรือซื้อกิจการที่ไหน",
         "อุปทาน–อุปสงค์–คุณภาพคู่แข่ง แยกรายรัฐ", "BQ1, BQ2, BQ8"],
        ["COO / ผู้จัดการภูมิภาค", "จะเข้าไปแทรกแซงที่ไหนก่อน",
         "รายชื่อสถานพยาบาลที่คะแนนกำลังตก เรียงลำดับความเร่งด่วน", "BQ6"],
        ["ฝ่ายบุคคล (HR)", "ระดับกำลังคน และงบรักษาคนไว้",
         "ความเชื่อมโยงระหว่างชั่วโมงพยาบาล อัตราลาออก และคุณภาพ", "BQ3, BQ4"],
        ["CFO / นักลงทุนสัมพันธ์", "ผลตอบแทนและความเสี่ยง",
         "ค่าปรับที่ต้องเผชิญ อัตราการเข้าพัก ตัวคูณจูงใจ", "BQ2, BQ5, BQ7"],
        ["ฝ่ายกำกับกฎระเบียบ", "ความเสี่ยงด้านกฎหมาย",
         "แนวโน้มค่าปรับ ข้อบกพร่องที่พบบ่อย สถานะ Special Focus", "BQ6, BQ7"],
    ]
    table(slide, rows, M, y, CONTENT_W, [2.3, 3.0, 4.4, 1.6],
          row_h=Inches(0.60), header_size=11, body_size=11)
    callout(slide, M, y + Inches(3.72), CONTENT_W, Inches(1.05), None,
            "**คอลัมน์กลางคือข้อสอบของกราฟทุกใบ** — ถ้ากราฟใบหนึ่งไม่ช่วยให้ใครสักคน"
            "ในตารางนี้ตัดสินใจข้อใดข้อหนึ่งได้ กราฟนั้นไม่ควรอยู่บนแดชบอร์ด "
            "เกณฑ์ข้อนี้คือเหตุผลที่แดชบอร์ดมี 13 กราฟ ไม่ใช่ 40 กราฟ",
            body_size=12)


def slide_questions(prs):
    slide = new_slide(prs)
    y = head(slide, "01 · BUSINESS QUESTIONS", "คำถามทางธุรกิจ 8 ข้อ (โจทย์กำหนดขั้นต่ำ 5 ข้อ)")
    rows = [
        ["#", "คำถาม", "ผู้ใช้", "ตอบจากตารางใด"],
        ["BQ1", "รัฐใดน่าเข้าไปลงทุนที่สุด เมื่อดูเตียงต่อผู้สูงอายุ 1,000 คน อัตราการเข้าพัก และคุณภาพคู่แข่ง",
         "CEO", "Fact 1 × Dim_Geography × Ref_State_Population"],
        ["BQ2", "รูปแบบการถือครองต่างกันอย่างไร ในด้านคุณภาพ กำลังคน และค่าปรับ",
         "CEO, CFO", "ทั้งสอง Fact × Dim_Ownership"],
        ["BQ3", "ชั่วโมงพยาบาลสัมพันธ์กับคะแนนและค่าปรับอย่างไร และจุดคุ้มอยู่ที่ใด",
         "HR, CFO", "Fact_Facility_Monthly"],
        ["BQ4", "อัตราการลาออกสูงถึงระดับใดจึงเริ่มฉุดคุณภาพและอัตราการเข้าพัก",
         "HR", "Fact_Facility_Monthly"],
        ["BQ5", "ปี 2562–2569 เกิดอะไรขึ้น ธุรกิจฟื้นจากโควิดแล้วหรือยัง",
         "CEO, CFO", "Fact 1 × Dim_Date"],
        ["BQ6", "สถานพยาบาลใดกำลังเสื่อมคุณภาพและมีแนวโน้มถูกลงโทษ",
         "COO, กำกับกฎ", "ทั้งสอง Fact × Dim_Facility (SCD2)"],
        ["BQ7", "ค่าปรับกระจุกตัวที่รัฐใดและช่วงเวลาใด รัฐใดบังคับใช้เข้มที่สุดต่อเตียง",
         "กำกับกฎ", "Fact 2 × Dim_Geography × Dim_Date"],
        ["BQ8", "เครือขนาดใหญ่ทำได้ดีกว่าจริง หรือแค่กดต้นทุนแรงงาน",
         "CEO", "Fact 1 × Dim_Chain"],
    ]
    table(slide, rows, M, y, CONTENT_W, [0.75, 6.6, 1.6, 4.0],
          row_h=Inches(0.46), header_size=11, body_size=10.5)
    callout(slide, M, y + Inches(4.30), CONTENT_W, Inches(0.86), None,
            "กำหนด 8 ข้อแทนที่จะเป็น 5 ข้อ เพราะคอลัมน์ขวาทำหน้าที่เป็น **ข้อสอบของ schema** ไปด้วย "
            "— คำถามข้อใดที่ไม่มีตารางรองรับ แปลว่า schema ยังออกแบบไม่เสร็จ",
            body_size=12)


def slide_measures(prs):
    slide = new_slide(prs)
    y = head(slide, "01 · MEASURES", "ตัววัดผลการดำเนินงาน 10 ตัว (โจทย์กำหนดขั้นต่ำ 3 ตัว)")
    rows = [
        ["Measure", "วิธีคำนวณ", "บวกได้ไหม", "ความหมายและประโยชน์"],
        ["M1 อัตราการเข้าพัก", "SUM(avg_residents_per_day) / SUM(certified_beds)",
         "ไม่ได้ (อัตราส่วน)", "ตัวขับรายได้หลัก ต่ำกว่า 80% มักแปลว่าขาดทุน"],
        ["M2 Resident-days", "avg_residents_per_day × จำนวนวันในงวด",
         "ได้เต็ม", "ปริมาณการดูแลที่ขายได้จริง เป็นน้ำหนักของ M1, M3, M4 และตัวหารของ M9"],
        ["M3 ชั่วโมงพยาบาล/คน/วัน", "reported_total_nurse_hprd ถ่วงด้วย M2",
         "ไม่ได้", "ตัวแทนต้นทุนแรงงาน และตัวทำนายคุณภาพที่แรงที่สุดตัวเดียว"],
        ["M4 อัตราลาออกพยาบาล", "total_nursing_turnover_pct ถ่วงด้วย M2",
         "ไม่ได้", "ต้นทุนแฝงด้านสรรหาและฝึกอบรม เป็นสัญญาณเตือนคุณภาพล่วงหน้า"],
        ["M5 มูลค่าค่าปรับ", "SUM(fine_amount_usd)",
         "ได้เต็ม", "ความเสี่ยงด้านกฎระเบียบที่แปลงเป็นตัวเงินแล้ว"],
        ["M6 จำนวนครั้งที่ถูกลงโทษ", "COUNT(*) บน Fact_Penalty_Event",
         "ได้เต็ม", "อ่านคู่กับ M5 เพื่อแยก “ปรับบ่อยแต่เบา” ออกจาก “ปรับครั้งเดียวหนัก”"],
        ["M7 คะแนนดาวเฉลี่ย", "AVG(overall_rating) สเกล 1–5",
         "ไม่ได้ (อันดับ)", "ตัวเลขที่ครอบครัวผู้ป่วยใช้เลือกจริง"],
        ["M8 จำนวนข้อบกพร่อง", "SUM(cycle1_total_deficiencies)",
         "ได้เต็ม", "สัญญาณเชิงปฏิบัติการที่มาก่อนค่าปรับ"],
        ["M9 ค่าปรับต่อ resident-day", "M5 / M2",
         "คำนวณตอน query", "ทำให้เทียบความเสี่ยงข้ามสถานพยาบาลคนละขนาดได้"],
        ["M10 เตียงต่อผู้สูงอายุ 1,000 คน", "SUM(certified_beds) / (pop_65plus / 1000)",
         "คำนวณตอน query", "ความอิ่มตัวของตลาด ค่าต่ำแปลว่าอุปสงค์ยังไม่ถูกตอบ (BQ1)"],
    ]
    table(slide, rows, M, y, CONTENT_W, [2.5, 3.9, 1.5, 5.0],
          row_h=Inches(0.40), header_size=10.5, body_size=9.8)
    callout(slide, M, y + Inches(4.52), CONTENT_W, Inches(0.62), None,
            "M1–M8 เก็บอยู่ในตาราง Fact · M9 และ M10 คำนวณตอน query · "
            "**ไม่มีอัตราส่วนตัวใดถูกเก็บเป็นคอลัมน์เลย โดยเจตนา** — เหตุผลอยู่ในสไลด์ถัดไป",
            body_size=11.5)


def slide_ratio_trap(prs):
    slide = new_slide(prs)
    y = head(slide, "01 · MEASURES", "กับดักที่พบบ่อยที่สุดในงานคลังข้อมูล: การเฉลี่ยอัตราส่วน",
             sub="คอลัมน์ “บวกได้ไหม” ในสไลด์ก่อนหน้าไม่ใช่รายละเอียดทางทฤษฎี "
                 "— มันคือสิ่งที่ตัดสินว่าแดชบอร์ดรวมยอดถูกหรือโกหกเงียบ ๆ")
    left_w = Inches(6.1)
    rect(slide, M, y, left_w, Inches(2.55), fill=WHITE, line=GRID)
    tf = textbox(slide, M + Inches(0.3), y + Inches(0.26), left_w - Inches(0.6), Inches(2.1))
    para(tf, "สถานพยาบาลสองแห่ง", size=13, color=INK, font=HEAD_FONT, bold=True,
         first=True, space_after=8)
    para(tf, "แห่งที่ 1 — 10 เตียง ผู้พักอาศัย 5 คน  →  **50%**", size=12.5, color=INK2,
         space_after=4)
    para(tf, "แห่งที่ 2 — 200 เตียง ผู้พักอาศัย 180 คน  →  **90%**", size=12.5, color=INK2,
         space_after=10)
    para(tf, "AVG(occupancy)  =  (50 + 90) / 2  =  **70%**   ✗ ผิด",
         size=13.5, color=WARN, space_after=4)
    para(tf, "SUM(residents) / SUM(beds)  =  185 / 210  =  **88.1%**   ✓ ถูก",
         size=13.5, color=GREEN, space_after=0)

    right_x = M + left_w + Inches(0.4)
    right_w = CONTENT_W - left_w - Inches(0.4)
    rect(slide, right_x, y, right_w, Inches(2.55), fill=WHITE, line=GRID)
    tf = textbox(slide, right_x + Inches(0.3), y + Inches(0.26), right_w - Inches(0.6),
                 Inches(2.1))
    para(tf, "และนี่ไม่ใช่ความกังวลในตำรา", size=13, color=INK, font=HEAD_FONT, bold=True,
         first=True, space_after=8)
    para(tf, "วัดจริงบนคลังนี้ ทุกแถวรวมแถวที่ติดธงสงสัย สองวิธีให้คำตอบต่างกันทุกงวด",
         size=11.5, color=INK2, space_after=10)
    table(slide, [
        ["งวด", "SUM/SUM (ถูก)", "AVG ของอัตราส่วน (ผิด)", "ส่วนต่าง"],
        ["2019-01", "78.98%", "80.69%", "+1.72 จุด"],
        ["2026-06", "79.45%", "80.27%", "+0.82 จุด"],
        ["2026-07", "80.24%", "81.25%", "+1.00 จุด"],
    ], right_x + Inches(0.3), y + Inches(1.06), right_w - Inches(0.6),
        [1.2, 1.6, 2.2, 1.2], row_h=Inches(0.31), header_size=10, body_size=9.8,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])

    callout(slide, M, y + Inches(2.78), CONTENT_W, Inches(2.1),
            "การป้องกันจึงเป็นเชิงโครงสร้าง ไม่ใช่การเตือนกันเอง",
            "ตาราง Fact **เก็บตัวตั้งกับตัวหารแยกกัน และไม่เก็บอัตราส่วนเป็นคอลัมน์เลย** "
            "เพราะการเก็บอัตราส่วนคือการเชิญชวนให้ใครสักคนเอาไปเฉลี่ย\n\n"
            "view ชื่อ v_check_ratio_additivity คำนวณการเปรียบเทียบนี้ใหม่ทุกครั้งที่ build "
            "และแดชบอร์ดแสดงผลของมันบนแท็บคุณภาพข้อมูล — "
            "ทั้งสองตัวเลขหน้าตาเหมือนอัตราการเข้าพักทั้งคู่ และตัวที่ผิดนั้นผิดมากพอที่จะ"
            "**สลับอันดับรัฐได้** ซึ่งเปลี่ยนคำตอบของ BQ1 ทั้งข้อ", body_size=12)


def slide_sources(prs):
    slide = new_slide(prs)
    y = head(slide, "02 · DATA SOURCES", "ห้าแหล่งข้อมูล — สามแหล่งดึงอัตโนมัติผ่าน HTTP")
    rows = [
        ["#", "แหล่งข้อมูล", "รูปแบบ", "ให้อะไรกับโครงงาน"],
        ["S1", "CMS Nursing Homes — snapshot ย้อนหลัง",
         "ZIP หลายแฟ้ม 88 งวด (ใช้ 32) แต่ละงวดมี 20 CSV",
         "แกนหลักของทุกอย่าง ข้อมูลผู้ให้บริการและค่าปรับทุกงวด — ทางเดียวที่จะได้ประวัติ"],
        ["S2", "CMS Provider Data Catalog REST API", "JSON ผ่าน HTTP",
         "ค้นว่ามีงวดใดบ้าง ขนาดและวันที่ และใช้ตรวจงวดล่าสุดหลัง load"],
        ["S3", "SNF Value-Based Purchasing", "CSV รายปีงบประมาณ โครงสร้างต่างกันแต่ละปี",
         "ตัวคูณจูงใจการจ่ายเงิน (extract แล้ว ยังไม่ได้ทำแบบจำลอง)"],
        ["S4", "ประชากรอายุ 65 ปีขึ้นไป — U.S. Census",
         "Nested JSON ผ่าน API ที่ต้องมีคีย์ พร้อมแฟ้ม CSV สำรอง",
         "ด้านอุปสงค์ของ BQ1 และเป็นตัวหารของ M10"],
        ["S5", "ตารางอ้างอิงของ CMS", "CSV",
         "คำอธิบายรหัสข้อบกพร่อง และค่าเฉลี่ยรายรัฐ"],
    ]
    table(slide, rows, M, y, CONTENT_W, [0.55, 3.2, 3.5, 5.6],
          row_h=Inches(0.52), header_size=10.5, body_size=10)
    tf = textbox(slide, M, y + Inches(3.28), CONTENT_W, Inches(0.32))
    para(tf, "โจทย์กำหนดอย่างน้อย 3 แหล่ง และอย่างน้อย 1 แหล่งต้องซับซ้อนกว่า CSV ตาราง — "
             "**ใช้ 5 แหล่ง** โดย S1 เป็นหลายแฟ้มหลายช่วงเวลา · S2 และ S4 เป็น API · "
             "และ S4 คืน nested JSON ที่ต้องประกอบจากตัวแปรช่วงอายุ 12 ตัวก่อนจึงจะมีความหมาย",
         size=11, color=INK2, first=True, space_after=0)

    traps = [
        ("ชื่อไฟล์มีอย่างน้อยสามแบบ",
         "ProviderInfo_Download.csv (2562) · NH_ProviderInfo_Jun2026.csv · "
         "4pq5-n9py_2026-07-29_NH_ProviderInfo_Jul2026.csv → glob ธรรมดาจับได้บางงวดและ"
         "พลาดงวดอื่นเงียบ ๆ จึงใช้ regex รายยุคเก็บไว้ที่เดียวใน config.FILE_PATTERNS"),
        ("บาง snapshot ไม่ครบ",
         "งวด 2026-08-06 ขนาด 3.5 MB มีเพียง 4 แฟ้ม และไม่มีแฟ้มค่าปรับเลย โค้ดที่สมมติว่า"
         "มี 6 แฟ้มจะพัง หรือแย่กว่านั้นคือโหลดค่าปรับ 0 แถวโดยไม่บ่น จึงตรวจรายการแฟ้ม"
         "ของทุกงวดก่อนประมวลผล และแฟ้มที่ขาดจะข้ามเฉพาะ Fact ที่พึ่งมัน"),
        ("โครงสร้างคอลัมน์เปลี่ยนตามยุค",
         "78 คอลัมน์ในปี 2562 เทียบกับ 99 คอลัมน์ในปี 2569 พร้อมการเปลี่ยนชื่อทั้งระบบ "
         "นี่คือปัญหาที่ใหญ่ที่สุดในแปดข้อ และถูกจัดการเป็นกฎ Q7 ในสไลด์ถัดไป"),
    ]
    cw = int((CONTENT_W - Inches(0.36)) / 3)
    ty = y + Inches(3.78)
    for i, (t, b) in enumerate(traps):
        x = M + i * (cw + Inches(0.18))
        rect(slide, x, ty, cw, Inches(1.42), fill=RGBColor(0xFD, 0xF1, 0xEC))
        rect(slide, x, ty, Inches(0.05), Inches(1.42), fill=ORANGE)
        tf = textbox(slide, x + Inches(0.24), ty + Inches(0.14), cw - Inches(0.42),
                     Inches(1.18))
        para(tf, "กับดักที่ %d · %s" % (i + 1, t), size=11, color=ORANGE, font=HEAD_FONT,
             bold=True, first=True, space_after=4)
        para(tf, b, size=9.3, color=INK2, space_after=0, line_spacing=1.14)


def slide_quality(prs):
    slide = new_slide(prs)
    y = head(slide, "02 · DATA QUALITY", "ปัญหาคุณภาพข้อมูล 8 ประเภท (โจทย์กำหนดขั้นต่ำ 3 ประเภท)")
    rows = [
        ["กฎ", "ประเภทปัญหา", "สิ่งที่พบจริงในข้อมูล", "วิธีจัดการ"],
        ["Q1", "Mismatched keys",
         "CCN เป็นรหัส 6 หลักที่มีศูนย์นำหน้า (015009) อ่านเป็นตัวเลขกลายเป็น 15009 แล้ว join หลุดโดยไม่มี error · ZIP ก็เป็นแบบเดียวกัน",
         "อ่านทุกคอลัมน์เป็น str ก่อน แล้วแปลงเฉพาะที่เป็นตัวเลขจริง"],
        ["Q2", "Duplicate records",
         "แฟ้มค่าปรับเป็นหน้าต่างหมุน 3 ปี ค่าปรับหนึ่งรายการจึงโผล่ซ้ำได้ถึง 36 งวด · Fine ID มีเฉพาะงวด 202606 ขึ้นไป",
         "สองระบบคีย์ตามยุค — fine_id 14,162 · natural key 65,641 · 88.3% ของแถวดิบเป็นของซ้ำ"],
        ["Q3", "Null และรหัสเชิงอรรถ",
         "ค่าที่ถูกปิดบังถูกเผยแพร่เป็นรหัสเชิงอรรถ ไม่ใช่ค่าว่าง จึงถูกเฉลี่ยราวกับเป็นข้อมูลจริง",
         "ตัดออกจากตัวหาร ไม่แทนด้วยศูนย์"],
        ["Q4", "Invalid values",
         "อัตราลาออกเกิน 100% (สูงสุด 1414%) · อัตราการเข้าพักเกิน 100% ก็มี แต่เป็นของจริงจากเตียงเสริม",
         "ติดธง ไม่ลบ — 1,207 แถว (0.25%) ติด is_suspect และแดชบอร์ดมีสวิตช์ให้เลือก"],
        ["Q5", "Category ไม่สอดคล้อง",
         "ชื่อสถานพยาบาลยุค 2562 พิมพ์ใหญ่เล็กไม่ตรงกัน ถ้าปล่อยไว้ราว 17% จะดูเหมือนเปลี่ยนชื่อ · เจ้าของมี 13 ประเภท",
         "ทำชื่อเป็นตัวพิมพ์ใหญ่ · ยุบ 13 ประเภทเหลือ 3 กลุ่ม (เปลี่ยนประเภทราวครึ่ง แต่เปลี่ยนกลุ่มเพียง ~10%)"],
        ["Q6", "รูปแบบวันที่ไม่ตรงกัน",
         "ทั้งสองปลายของช่วงใช้ YYYY-MM-DD ยังไม่พบความต่าง แต่มีอีกหลายงวดที่ยังไม่ได้เปิดดูทีละคอลัมน์",
         "คง parser หลายรูปแบบไว้ · ให้คะแนนแต่ละรูปแบบกับทั้งคอลัมน์ ไม่ fallback ทีละแถว · ความกำกวมถูกบันทึก ไม่เดา"],
        ["Q7", "Schema drift ข้ามปี",
         "ใหญ่ที่สุดในแปดข้อ — 78 คอลัมน์ปี 2562 เทียบ 99 คอลัมน์ปี 2569 พร้อมการเปลี่ยนชื่อทั้งระบบ",
         "ตารางแม็ปคอลัมน์รายยุคใน columns.py อ้างอิงจากพจนานุกรมข้อมูลของ CMS โดยตรง"],
        ["Q8", "หน่วยวัดไม่ตรงกัน",
         "ยุค 2569 เก็บอัตราลาออกเป็นเปอร์เซ็นต์ · ยุค 2562 ไม่มีคอลัมน์นี้เลย จึงเทียบข้ามยุคตรง ๆ ไม่ได้",
         "ปรับหน่วยให้ตรงกันเท่าที่เทียบได้ · และประกาศว่า BQ4 เป็น cross-sectional ไม่ใช่ longitudinal"],
    ]
    table(slide, rows, M, y, CONTENT_W, [0.5, 1.9, 5.6, 4.0],
          row_h=Inches(0.455), header_size=10.5, body_size=9.2)

    ty = y + Inches(4.20)
    cw = int((CONTENT_W - Inches(0.24)) / 2)
    for i, (t, b, col, bg) in enumerate([
        ("นโยบายที่ 1 · ติดธง ไม่ลบ",
         "แถวที่ไม่ผ่านกฎถูกทำเครื่องหมาย is_suspect แล้วเก็บไว้ ไม่ทิ้ง การลบทิ้งคือการซ่อน"
         "ความบกพร่องจากคนที่จำเป็นต้องเห็นมันที่สุด และค่าผิดปกติบางตัวเป็นของจริง "
         "— แดชบอร์ดจึงเปิดสวิตช์ให้ผู้อ่านเลือกเอง", ACCENT, ACCENT_TINT),
        ("นโยบายที่ 2 · ทุกกฎเขียนลงตาราง ไม่ใช่ console",
         "กฎทั้งแปดข้อ append ลง etl_run_log ในคลัง — ตอนนี้ **1,555 แถว จากกฎ 20 ชนิด** "
         "แท็บคุณภาพข้อมูลบนแดชบอร์ดอ่านตารางนั้นสด ๆ สิ่งที่รายงานอ้างกับสิ่งที่แดชบอร์ดแสดง"
         "จึงหลุดจากกันไม่ได้", GREEN, RGBColor(0xEA, 0xF7, 0xF1)),
    ]):
        x = M + i * (cw + Inches(0.24))
        rect(slide, x, ty, cw, Inches(1.0), fill=bg)
        rect(slide, x, ty, Inches(0.05), Inches(1.0), fill=col)
        tf = textbox(slide, x + Inches(0.24), ty + Inches(0.12), cw - Inches(0.42),
                     Inches(0.78))
        para(tf, t, size=11, color=col, font=HEAD_FONT, bold=True, first=True, space_after=3)
        para(tf, b, size=9.6, color=INK2, space_after=0, line_spacing=1.12)


def slide_star(prs):
    slide = new_slide(prs)
    y = head(slide, "02 · DATA WAREHOUSE DESIGN",
             "Star schema — 2 Fact Table และ 6 Dimension Table",
             sub="ออกแบบตามสี่ขั้นของ Kimball · Dim ทั้งห้าตัวที่ขนาบอยู่เป็น conformed dimension "
                 "ที่ Fact ทั้งสองใช้ร่วมกัน · Dim_Penalty_Type ใช้เฉพาะ Fact 2")
    star_schema(slide, y - Inches(0.06))


def slide_grain(prs):
    slide = new_slide(prs)
    y = head(slide, "02 · GRAIN", "หนึ่งแถวของ Fact Table แปลว่าอะไร")
    rows = [
        ["Fact Table", "หนึ่งแถวหมายถึง", "ชนิด"],
        ["Fact_Facility_Monthly",
         "สถานพยาบาลหนึ่งแห่ง (หนึ่ง CCN) ในสภาพที่เป็นอยู่ ณ วันที่ CMS ประมวลผลหนึ่งงวด",
         "Periodic snapshot — มีแถวทุกงวดแม้ไม่มีอะไรเปลี่ยน"],
        ["Fact_Penalty_Event",
         "การลงโทษหนึ่งครั้ง (ค่าปรับ หรือการระงับการจ่ายเงิน) ต่อสถานพยาบาลหนึ่งแห่งในวันหนึ่ง",
         "Transaction — ไม่มีเหตุการณ์ ก็ไม่มีแถว"],
    ]
    table(slide, rows, M, y, CONTENT_W, [2.9, 6.3, 4.2],
          row_h=Inches(0.62), header_size=11, body_size=11)
    tf = textbox(slide, M, y + Inches(2.02), CONTENT_W, Inches(0.3))
    para(tf, "การมีสอง Fact ได้คะแนนพิเศษตามที่โจทย์เสนอไว้ แต่เหตุผลที่มีสองตัวคือ **grain ต่างกัน** ไม่ใช่คะแนน",
         size=11.5, color=MUTED, first=True, space_after=0)

    cw = int((CONTENT_W - Inches(0.3)) / 2)
    cy = y + Inches(2.45)
    callout(slide, M, cy, cw, Inches(2.55),
            "ทำไมสอง Fact นี้ห้ามรวมกัน",
            "ทั้งคู่ตอบคำถามที่ grain ต่างกัน ถ้า JOIN ตรง ๆ ทีละแถว ฝั่ง snapshot จะถูกคูณ"
            "ด้วยจำนวนครั้งที่ถูกลงโทษ\n\n"
            "วัดจริงบนคลังนี้ — เตียงที่ผ่านการรับรองของกลุ่ม for-profit พองจาก "
            "**1,203,047 เป็น 2,586,212 คือ 2.1 เท่า โดยไม่มี error ใด ๆ**\n\n"
            "ทุก measure ที่ข้าม Fact จึงรวมแต่ละฝั่งลงมาที่ระดับเดียวกันก่อน แล้วค่อยเชื่อมยอด "
            "เช่น view v_fines_per_resident_day ที่ลดทั้งสองฝั่งเหลือ (รัฐ, ปี) ก่อน FULL JOIN",
            accent=WARN, fill=RGBColor(0xFB, 0xEE, 0xEB))
    callout(slide, M + cw + Inches(0.3), cy, cw, Inches(2.55),
            "ทำไม PK ของ Fact 1 คือ (snapshot_date_key, ccn) ไม่ใช่ facility_key",
            "facility_key เป็น surrogate แบบ SCD2 สถานพยาบาลหนึ่งแห่งจึงเป็นเจ้าของค่านี้"
            "หลายค่า\n\n"
            "คีย์ที่สร้างบน facility_key จะยอมรับ **สองแถวของสถานพยาบาลเดียวกันในงวดเดียวกัน**"
            "ได้อย่างถูกกฎ ตราบใดที่มีเส้นแบ่งเวอร์ชันคั่นอยู่ตรงกลาง — ซึ่งคือแถวซ้ำที่ grain "
            "ห้ามไว้พอดี\n\n"
            "ccn คือตัวตนที่รอดจากการทำเวอร์ชัน จึงถูกเก็บไว้ในตาราง Fact เองในฐานะ "
            "degenerate dimension และใช้เป็นส่วนหนึ่งของคีย์")


def slide_built(prs):
    slide = new_slide(prs)
    y = head(slide, "02 · WHAT WAS BUILT", "สิ่งที่สร้างได้จริง — อ่านออกมาจาก eldercare.duckdb")
    rows = [
        ["ตาราง", "จำนวนแถว", "หมายเหตุ"],
        ["Dim_Date", "24,838", "1960-01-01 ถึง 2027-12-31 พร้อมแถว date_key = -1"],
        ["Dim_Facility", "55,425", "16,155 สถานพยาบาล เก็บประวัติแบบ SCD2 · เวอร์ชันปัจจุบัน 16,154"],
        ["Dim_Geography", "9,747", "grain (ZIP, เมือง, รัฐ) · ครอบคลุม 54 รัฐและดินแดน"],
        ["Dim_Ownership", "14", "13 ประเภท ยุบเหลือ 3 กลุ่ม"],
        ["Dim_Chain", "714", "712 เครือ บวก Independent (ไม่สังกัดเครือ) และ Unknown"],
        ["Dim_Penalty_Type", "3", "Fine · Payment Denial · Unknown"],
        ["Fact_Facility_Monthly", "483,183", "32 งวด (2562-01 → 2569-07) · 1,207 แถวติดธง is_suspect"],
        ["Fact_Penalty_Event", "79,803", "ค่าปรับ 70,566 ครั้ง รวม $1,616,747,274 · ระงับการจ่ายเงิน 9,237 ครั้ง 280,604 วัน"],
        ["Ref_State_Population", "255", "grain (รัฐ, ปี) — เจตนาไม่เก็บไว้บน Dim_Geography"],
        ["etl_run_log", "1,555", "ผลของกฎคุณภาพทุกข้อ · 20 ชนิดกฎ"],
    ]
    table(slide, rows, M, y, Inches(7.55), [2.5, 1.3, 5.6],
          row_h=Inches(0.375), header_size=10.5, body_size=9.6,
          aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.LEFT])

    x = M + Inches(7.85)
    w = CONTENT_W - Inches(7.85)
    callout(slide, x, y, w, Inches(1.58),
            "ทุก Dimension มีแถว key = -1 “Unknown”",
            "เพื่อไม่ให้ Fact ถือ foreign key ที่เป็น null — null จะทำให้แถวหายเงียบ ๆ ตอน join · "
            "Dim_Chain แยก “ไม่สังกัดเครือ” (key = 0) ออกจาก “ไม่ทราบ” (−1) เพราะปี 2562 "
            "ไม่มีคอลัมน์เครือเลย", body_size=10)
    callout(slide, x, y + Inches(1.74), w, Inches(2.10),
            "Slowly Changing Dimension แบบที่ 2",
            "Dim_Facility เก็บประวัติด้วย effective_date / expiry_date / is_current · "
            "เพราะ facility_key ไม่ได้เป็นหนึ่งต่อหนึ่งกับ ccn ทุก join จาก Fact มาที่ Dimension นี้"
            "ต้องเลือกเวอร์ชันที่ช่วงเวลาครอบคลุมวันที่ของ fact — กฎนี้อยู่ในฟังก์ชันเดียวคือ "
            "facts.lookup_facility_key() และไม่เคยถูกเขียนซ้ำแบบ inline",
            accent=GREEN, fill=RGBColor(0xEA, 0xF7, 0xF1), body_size=10)
    callout(slide, x, y + Inches(4.00), w, Inches(1.24),
            "Semantic layer — 10 views",
            "views.sql นิยาม M1–M10 ไว้ที่เดียว แดชบอร์ด query view ไม่ใช่ตาราง Fact "
            "เพราะสามในสิบตัวจะให้คำตอบที่ผิดแบบดูน่าเชื่อถ้าคำนวณตรง ๆ", body_size=10)


def slide_etl(prs):
    slide = new_slide(prs)
    y = head(slide, "02 · ETL PROCESS", "Extract → Clean → Transform → Integrate → Load",
             sub="ทั้งกระบวนการเป็นคำสั่งห้าบรรทัด และไม่ต้องแก้ข้อมูลด้วยมือในขั้นตอนใดเลย")
    etl_diagram(slide, y + Inches(0.05))
    ty = y + Inches(2.52)
    cw = int((CONTENT_W - Inches(0.3)) / 2)
    callout(slide, M, ty, cw, Inches(1.52),
            "แต่ละแหล่งถูกนำเข้าอย่างไร",
            "fetch_snapshots.py เรียก archive API ซึ่งคืน 96 รายการ — 88 รายการเป็นงวดจริง "
            "และ 8 รายการเป็นชุดรายปีที่ซ้ำกับงวดจริง เก็บเฉพาะแบบแรก · "
            "ดาวน์โหลดแบบ incremental (งวดที่มีบนดิสก์แล้วจะถูกข้าม) และแตกเฉพาะ 6 CSV ที่ใช้จริง · "
            "ประชากรมาจาก population.py ที่เลือก Census ACS API เมื่อมี CENSUS_API_KEY "
            "และตกไปใช้แฟ้มสำรองเมื่อไม่มี — ทั้งสองทางลงเอยที่รูปร่างเดียวกัน",
            body_size=10.2)
    callout(slide, M + cw + Inches(0.3), ty, cw, Inches(1.52),
            "หลายแหล่งถูกเชื่อมกันอย่างไร",
            "ccn เป็นคีย์เชื่อมตลอดสาย ซึ่งเป็นเหตุผลที่กฎศูนย์นำหน้า (Q1) ต้องบังคับตั้งแต่ตอนอ่าน "
            "ไม่ใช่มาแก้ทีหลัง · Fact เชื่อมกับ Dim_Facility ผ่านการค้นหาแบบ SCD2 · "
            "งวด 2026-07-29 และ 2026-08-06 ใช้วันประมวลผลเดียวกัน ถ้าไม่กระทบยอด PK ของ Fact 1 "
            "จะชนกันทันที · ประชากรเชื่อมที่ระดับ (รัฐ, ปี) และเจตนาไม่เขียนลง Dim_Geography "
            "เพราะตารางนั้น grain เป็น ZIP — SUM เผลอ ๆ จะคูณ Texas ด้วยจำนวน ZIP",
            body_size=10.2)
    tf = textbox(slide, M, ty + Inches(1.66), CONTENT_W, Inches(0.4))
    para(tf, "**เครื่องมือ** Python 3 + pandas สำหรับการแปลงข้อมูลทั้งหมด · DuckDB เป็นคลังแบบไฟล์เดียว · "
             "Streamlit + matplotlib/Plotly สำหรับแดชบอร์ด · "
             "**รันซ้ำได้** dimension ถูกสร้างใหม่ทั้งหมดทุกครั้ง ไม่มีสถานะค้างข้ามรอบ และสาม build ติดกันให้ผลเหมือนกัน",
         size=11, color=INK2, first=True, space_after=0)


def slide_verification(prs):
    slide = new_slide(prs)
    y = head(slide, "02 · VERIFICATION", "ตรวจสอบความถูกต้องสามระดับ: ก่อน · ระหว่าง · หลัง Load")
    levels = [
        ("ก่อน Load", "load.py", ACCENT,
         "ตรวจ primary key, foreign key และช่วงเวลาของ SCD2 บน DataFrame ที่ยังอยู่ในหน่วยความจำ "
         "ปัญหาถูกจับก่อนที่จะมีอะไรถูกเขียนลงคลัง"),
        ("ระหว่าง Load", "build_warehouse.py", ORANGE,
         "คัดลอกแต่ละตารางออกไป สร้างใหม่จาก schema.sql พร้อม key / uniqueness / CHECK "
         "แล้วใส่ข้อมูลกลับผ่านข้อจำกัดเหล่านั้น — **การ insert เองคือการทดสอบ** ไม่ใช่คำประกาศว่าทดสอบแล้ว "
         "แถวเสียทำให้ run ล้มและคืนสถานะเดิม (ทดสอบแล้วว่ากัดจริง: facility_key = 88888888 "
         "ถูกปฏิเสธเป็น FK violation และ overall_rating = 9 เป็น CHECK violation)"),
        ("หลัง Load", "ชุดตรวจ 4 ชุด", GREEN,
         "query คลังด้วย SQL ที่เขียนแยกจากโค้ดที่สร้างข้อมูล — ถ้าใช้โค้ดเดียวกันตรวจตัวเอง "
         "ความผิดพลาดเดียวกันจะผ่านทั้งสองฝั่ง"),
    ]
    cw = int((CONTENT_W - Inches(0.36)) / 3)
    for i, (name, mod, col, body) in enumerate(levels):
        x = M + i * (cw + Inches(0.18))
        rect(slide, x, y, cw, Inches(2.05), fill=WHITE, line=GRID)
        rect(slide, x, y, cw, Inches(0.05), fill=col)
        tf = textbox(slide, x + Inches(0.24), y + Inches(0.22), cw - Inches(0.48), Inches(1.7))
        para(tf, name, size=15, color=col, font=HEAD_FONT, bold=True, first=True,
             space_after=1, line_spacing=1.0)
        para(tf, mod, size=9.5, color=MUTED, space_after=6)
        para(tf, body, size=10, color=INK2, space_after=0, line_spacing=1.14)

    ty = y + Inches(2.26)
    table(slide, [
        ["ชุดตรวจ", "ตรวจอะไร", "ผล"],
        ["verify_dims.py", "Dimension ทั้งหก รวมความสมบูรณ์ของ SCD2", "24 / 24"],
        ["verify_facts.py", "Fact ทั้งสอง รวมการกระทบยอดกับ CMS และการอ่านครบทุกแฟ้ม", "24 / 24"],
        ["verify_dashboard.py", "measure บนแดชบอร์ดเทียบกับ view ในคลัง ที่ความละเอียด 1e-9", "56 / 56"],
        ["test_dashboard.py", "รันแอปจริงภายใต้ชุดตัวกรอง 11 แบบ", "11 / 11"],
    ], M, ty, Inches(6.35), [2.0, 5.0, 1.2], row_h=Inches(0.42),
        header_size=10.5, body_size=10,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT])

    x = M + Inches(6.65)
    callout(slide, x, ty, CONTENT_W - Inches(6.65), Inches(2.10),
            "หลักฐานที่แข็งที่สุดไม่ใช่ชุดตรวจเหล่านี้",
            "ชุดตรวจที่เขียนโดยคนกลุ่มเดียวกับ pipeline อาจเห็นด้วยกันแล้วผิดทั้งคู่ได้ "
            "หลักฐานที่แข็งกว่าคือการกระทบยอดกับตัวเลขที่ **CMS คำนวณขึ้นเองอย่างเป็นอิสระ**\n\n"
            "จำกัด Fact_Penalty_Event ให้เหลือเฉพาะค่าปรับที่อยู่ในแฟ้มงวด มิ.ย. 2569 "
            "แล้วรวมยอดรายสถานพยาบาล ได้ตรงกับคอลัมน์ Total Amount of Fines in Dollars ของ CMS "
            "**ครบ 6,563 จาก 6,563 แห่ง**",
            accent=GREEN, fill=RGBColor(0xEA, 0xF7, 0xF1))
    tf = textbox(slide, M, ty + Inches(2.28), CONTENT_W, Inches(0.5))
    para(tf, "**การกำจัดของซ้ำก็ถูกวัดเช่นกัน** — 32 แฟ้มค่าปรับที่แตกออกมา ถูกอ่านครบทั้ง 32 แฟ้ม "
             "รวม 680,900 แถวดิบ เหลือ 79,803 เหตุการณ์ แปลว่า **88.3% ของแถวดิบเป็นของซ้ำ** "
             "ซึ่งเป็นสิ่งที่หน้าต่างหมุน 3 ปีต้องให้ผลออกมาแบบนั้นพอดี",
         size=11, color=INK2, first=True, space_after=0)


def slide_dashboard(prs):
    slide = new_slide(prs)
    y = head(slide, "03 · DASHBOARD", "แดชบอร์ด Streamlit ที่อ่านคลังแบบ read-only")
    table(slide, [
        ["ข้อกำหนดของโจทย์", "ขั้นต่ำ", "ทำได้"],
        ["Measures สรุปผลประกอบการ", "≥ 3", "10"],
        ["กราฟ", "≥ 5", "13"],
        ["การวิเคราะห์ตามช่วงเวลา", "≥ 1", "2"],
        ["การเปรียบเทียบข้อมูล", "≥ 1", "2"],
        ["Filter / Interactive control", "≥ 2", "5"],
        ["Insight และข้อเสนอแนะ", "≥ 5", "8"],
    ], M, y, Inches(4.6), [3.0, 1.0, 1.0], row_h=Inches(0.40),
        header_size=10.5, body_size=10.5,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
    tf = textbox(slide, M, y + Inches(3.02), Inches(4.6), Inches(0.5))
    para(tf, "verify_dashboard.py **นับตัวเลขคอลัมน์ขวาจากโค้ดจริง** ไม่ใช่จากความจำ "
             "ตารางนี้จึงล้าสมัยเงียบ ๆ ไม่ได้",
         size=10, color=MUTED, first=True, space_after=0)

    tf = textbox(slide, M, y + Inches(3.62), Inches(4.6), Inches(1.6))
    para(tf, "แท็บทั้งเจ็ด", size=12, color=INK, font=HEAD_FONT, bold=True, first=True,
         space_after=5)
    para(tf, "ภาพรวม · ตลาด (BQ1) · ผู้ประกอบการ (BQ2, BQ8) · กำลังคน (BQ3, BQ4) · "
             "แนวโน้ม (BQ5) · ความเสี่ยง (BQ6, BQ7) · ข้อเสนอแนะ  "
             "และแท็บคุณภาพข้อมูลที่อ่าน etl_run_log สด ๆ",
         size=10.5, color=INK2, space_after=6)
    para(tf, "**Insight ถูกคำนวณตอน render ไม่ได้เขียนฝังไว้ในหน้า** ข้อความกับกราฟอ่านตัวเลข"
             "ชุดเดียวกันใน request เดียวกัน จึงขัดกันเองไม่ได้ และการขยับตัวกรองจะคำนวณ"
             "ข้อสรุปใหม่ทั้งหมด",
         size=10.5, color=INK2, space_after=0)

    place_image(slide, SHOTS / "01_overview.png", M + Inches(4.95), y,
                CONTENT_W - Inches(4.95), Inches(5.16))


def slide_filters(prs):
    slide = new_slide(prs)
    y = head(slide, "03 · INTERACTIVE CONTROLS", "ตัวกรองห้ารายการ — ทุกกราฟและทุก measure คำนวณใหม่")
    w = int((CONTENT_W - Inches(0.3)) / 2)
    place_image(slide, SHOTS / "08_filtered_IL_forprofit.png", M, y, w, Inches(3.4))
    place_image(slide, SHOTS / "09_filtered_watchlist.png", M + w + Inches(0.3), y, w,
                Inches(3.4))
    for i, cap in enumerate([
        "กรองรัฐ = IL และกลุ่มเจ้าของ = For profit — ตัวเลขบนแถบ measure ขยับทั้งแถบ",
        "แท็บความเสี่ยง (BQ6) หลังกรองเหลือเฉพาะรายชื่อเฝ้าระวัง พร้อมปุ่มดาวน์โหลดเป็น CSV",
    ]):
        tf = textbox(slide, M + i * (w + Inches(0.3)), y + Inches(3.5), w, Inches(0.4))
        para(tf, cap, size=10.5, color=INK2, first=True, space_after=0)

    ty = y + Inches(4.0)
    cw = int((CONTENT_W - Inches(0.3)) / 2)
    callout(slide, M, ty, cw, Inches(1.16), "ตัวควบคุมที่มี",
            "งวดข้อมูล · รัฐ · กลุ่มเจ้าของ · ขนาดเครือ · สวิตช์รวม/ไม่รวมแถวที่ติดธงสงสัย "
            "และมีตัวควบคุมเฉพาะแท็บอีกห้ารายการ เช่น เกณฑ์คะแนนเสี่ยงและจำนวนรัฐที่แสดง")
    callout(slide, M + cw + Inches(0.3), ty, cw, Inches(1.16),
            "ภาพนี้อ้างการโต้ตอบที่ไม่เคยเกิดไม่ได้",
            "make_screenshots.py อ่านค่า measure หัวตาราง **ก่อนและหลัง** ใส่ตัวกรอง "
            "แล้วปฏิเสธที่จะบันทึกภาพถ้าค่านั้นไม่ขยับ",
            accent=GREEN, fill=RGBColor(0xEA, 0xF7, 0xF1))


def slide_bq(prs, kicker, title, answer, points, reco, image, *, reco_h=Inches(1.05)):
    """เลย์เอาต์มาตรฐานของสไลด์คำตอบ: ข้อความซ้าย กราฟขวา ข้อเสนอแนะปิดท้าย"""
    slide = new_slide(prs)
    y = head(slide, kicker, title)
    lw = Inches(5.15)
    tf = textbox(slide, M, y, lw, Inches(0.3))
    _rich(tf.paragraphs[0], "คำตอบ", size=10, color=MUTED, font=HEAD_FONT, bold=True)
    tf = textbox(slide, M, y + Inches(0.26), lw, Inches(0.9))
    para(tf, answer, size=15.5, color=ACCENT, font=HEAD_FONT, bold=True, first=True,
         space_after=0, line_spacing=1.15)

    tf = textbox(slide, M, y + Inches(1.18), lw, Inches(3.0))
    bullets(tf, points, size=11.5, gap=10, first=True)

    reco_top = BODY_BOTTOM - reco_h
    callout(slide, M, reco_top, lw, reco_h, "ข้อเสนอแนะ", reco,
            accent=GREEN, fill=RGBColor(0xEA, 0xF7, 0xF1), title_size=11, body_size=11)

    ix = M + lw + Inches(0.3)
    place_image(slide, DASH / image, ix, y, CONTENT_W - lw - Inches(0.3), BODY_H)


def slide_bq1(prs):
    slide_bq(
        prs, "03 · BQ1 · ตลาด", "รัฐใดน่าเข้าไปลงทุนที่สุด",
        "New Mexico · Virginia · South Carolina",
        ["สามรัฐนี้ครบทั้งสามเงื่อนไขพร้อมกัน — อุปทานต่ำกว่าค่าเฉลี่ยประเทศ "
         "อัตราการเข้าพักสูงกว่าค่าเฉลี่ย และคะแนนคู่แข่งต่ำกว่า 3.00 ดาว",
         "New Mexico: **16.1 เตียง** ต่อผู้สูงอายุ 1,000 คน · เข้าพัก **83.4%** · คู่แข่ง **2.85 ดาว** "
         "เทียบระดับประเทศ 25.6 เตียง · 80.2% · 2.99 ดาว",
         "ควรเลี่ยง **Arkansas** (42.9 เตียง เข้าพัก 70.3%) · **Indiana** (40.3 · 75.2%) · "
         "**Iowa** (41.1 · 80.3%) — ตลาดล้นแล้ว การเพิ่มเตียงคือการแย่งผู้พักอาศัยจากบ้านที่มีอยู่",
         "Guam และ Puerto Rico ไม่ได้คะแนน เพราะไม่มีข้อมูลประชากร 65+ — "
         "ปล่อยว่างไว้ดีกว่าใส่ค่าประมาณแล้วให้ติดอันดับปลอม"],
        "ชอร์ตลิสต์สามรัฐแรกไว้ก่อน และอ่านคู่กับ BQ5 เสมอ — การตัดสินใจเรื่องกำลังการผลิต"
        "ต้องดูอุปสงค์**รายรัฐ** ไม่ใช่ตัวเลขรวมของทั้งประเทศ",
        "bq1_market_scatter.png")


def slide_bq2_bq8(prs):
    slide_bq(
        prs, "03 · BQ2 + BQ8 · ผู้ประกอบการ",
        "รูปแบบการถือครองและขนาดเครือให้ผลต่างกันอย่างไร",
        "For-profit ต่ำกว่าทุกมิติคุณภาพ · ขนาดเครือซื้ออัตราการเข้าพัก ไม่ได้ซื้อคุณภาพ",
        ["**BQ2** ดาวเฉลี่ย 2.80 / 3.57 / 3.28 · ชม.พยาบาล 3.62 / 4.21 / 4.17 · "
         "ค่าปรับต่อเตียง $337 / $216 / $228 (for-profit / non-profit / government)",
         "ค่าปรับต่อเตียงของ for-profit สูงกว่า non-profit **1.56 เท่า** และช่องว่างดาว 0.77 "
         "เดินคู่ไปกับชั่วโมงพยาบาลที่น้อยกว่า **16%** เสมอ ไม่ได้แยกจากกัน",
         "**BQ8** อิสระ 4.09 ชม. เทียบเครือ 50+ แห่ง 3.53 ชม. (**−13.8%**) ดาวต่ำกว่า **0.32** "
         "ลาออกสูงกว่า 4.6 จุด — แต่อัตราการเข้าพักสูงกว่า 1.6 จุด",
         "คุมตัวแปรเจ้าของแล้วยังจริง: ชั่วโมงพยาบาลลดลง**ทุกขั้น**ตามขนาดเครือ ในทั้งสามกลุ่ม "
         "(for-profit 3.75→3.54 · non-profit 4.38→3.47 · government 4.56→3.34) "
         "ขนาดเครือจึงมีผลที่รูปแบบการถือครองอธิบายไม่ได้"],
        "ในดีลซื้อกิจการ ให้ตั้งราคาความเสี่ยงของ for-profit เป็นปัจจัยที่ต้องจ่าย ไม่ใช่คุณลักษณะที่เป็นกลาง "
        "· อย่าจ่ายพรีเมียมด้าน**คุณภาพ**ให้ขนาดเครือ ให้จ่ายตามข้อได้เปรียบด้านอัตราการเข้าพัก "
        "แล้วกันงบชั่วโมงพยาบาลไว้ต่างหาก",
        "bq2_segments.png", reco_h=Inches(1.25))


def slide_bq3(prs):
    slide_bq(
        prs, "03 · BQ3 · กำลังคน", "ชั่วโมงพยาบาลกับคุณภาพและค่าปรับ จุดคุ้มอยู่ที่ใด",
        "จุดคุ้มอยู่ที่ 3.48 ชม./คน/วัน ซึ่งตรงกับเกณฑ์ขั้นต่ำของ CMS พอดี",
        ["ต่ำกว่า 3.48 ชม. ค่าปรับต่อเตียง **$404** เทียบกับ **$268** เมื่อถึงเกณฑ์ — สูงกว่า "
         "**1.51 เท่า** · และเคยถูกปรับ 52.6% เทียบกับ 39.7%",
         "**มีสองจุดคุ้ม ไม่ใช่จุดเดียว** เหนือราว 3.75 ชม. เส้นค่าปรับแบนแล้ว "
         "($289 → $264 → $283) ชั่วโมงที่เพิ่มยังซื้อดาวได้ แต่ไม่ซื้อการลดความเสี่ยงอีกต่อไป",
         "**กับดักที่ปิดไปแล้ว** CMS เอาชั่วโมงพยาบาลไปคิดเป็นดาวด้านการจัดพยาบาล ซึ่งเป็น"
         "องค์ประกอบของคะแนนรวมอยู่แล้ว การพล็อตชั่วโมงกับคะแนนรวมจึงเป็นวงกลมบางส่วน "
         "→ ข้อสรุปทั้งหมดจึงยืนบน **ดาวการตรวจสุขภาพ** (ผู้ตรวจให้คะแนนหน้างาน ไม่มีชั่วโมงในสูตร) "
         "และ **ค่าปรับ** แทน",
         "ความชันยังอยู่ครบเมื่อคุมตัวแปรกลุ่มเจ้าของทั้งสามกลุ่ม"],
        "**5,134 แห่ง (36%)** ยังอยู่ต่ำกว่าเกณฑ์ — งบกำลังคนก้อนถัดไปเป็นของพวกเขา "
        "ไม่ใช่การเติมชั่วโมงให้บ้านที่ผ่านเกณฑ์ไปแล้ว",
        "bq3_staffing_gradient.png")


def slide_bq4(prs):
    slide_bq(
        prs, "03 · BQ4 · กำลังคน", "อัตราการลาออกสูงแค่ไหนจึงเริ่มฉุดธุรกิจ",
        "มีสองเส้น และคนละคนใช้คนละเส้น",
        ["**คุณภาพไม่มีระดับปลอดภัย** คะแนนดาวลดตั้งแต่แบนด์แรกสุดและลดต่อเนื่อง "
         "ทุก 5 จุดของอัตราลาออกแลกไปราว 0.2 ดาว",
         "**อัตราการเข้าพักหักที่ 40%** ต่ำกว่านั้นทรงตัวราว 83.5% แล้วเสีย **1.57 จุด**ทันที"
         "ในช่วง 40–45% และพังที่ 65% ขึ้นไป เหลือ 71.7% (ค่าเฉลี่ยประเทศอยู่ที่ 44.7%)",
         "**การลาออกไม่ใช่แค่เงาของการจัดพยาบาลบาง** ในควอไทล์ที่จัดพยาบาลดีที่สุด "
         "การขยับจากลาออก <40% ไป 60%+ ยังทำคะแนนตก **1.34 ดาว** (4.15 → 2.81) "
         "ซึ่ง*มากกว่า*ช่วงห่างเดียวกันในควอไทล์ที่บางที่สุด → เป็นคานงัดอิสระ",
         "ข้อจำกัด: คอลัมน์อัตราลาออกมีเฉพาะยุค 2569 ข้อสรุปข้อนี้จึงเป็น cross-sectional"],
        "ฝ่ายบุคคลควรถือว่าอัตราลาออกทุกระดับเหนือศูนย์คุ้มที่จะแก้ ส่วน CFO ควรป้องเส้น **40%** "
        "เพราะนั่นคือจุดที่เงินอยู่ · โปรแกรมรักษาคนถูกกว่าการจ้างเพิ่มถาวร และหลักฐานบอกว่าได้ผลอย่างอิสระ",
        "bq4_turnover_gradient.png", reco_h=Inches(1.15))


def slide_bq5(prs):
    slide = new_slide(prs)
    y = head(slide, "03 · BQ5 · แนวโน้ม", "ปี 2562–2569 ธุรกิจฟื้นจากโควิดแล้วหรือยัง")
    lw = Inches(5.6)
    tf = textbox(slide, M, y, lw, Inches(0.3))
    _rich(tf.paragraphs[0], "คำตอบ", size=10, color=MUTED, font=HEAD_FONT, bold=True)
    tf = textbox(slide, M, y + Inches(0.26), lw, Inches(0.9))
    para(tf, "ฟื้นเต็มด้านการเข้าพักแล้ว แต่ใช้เวลาห้าปี — ส่วนคุณภาพยังไม่ฟื้น",
         size=15.5, color=ACCENT, font=HEAD_FONT, bold=True, first=True, space_after=0,
         line_spacing=1.15)
    table(slide, [
        ["", "ก่อนโควิด ก.ค. 62", "ก้นเหว ก.ค. 64", "ปัจจุบัน ก.ค. 69"],
        ["อัตราการเข้าพัก", "80.7%", "67.4%", "80.2%"],
        ["ผู้พักอาศัย", "1,327,186", "1,096,242", "1,256,235"],
        ["ชม.พยาบาล/คน/วัน", "3.74", "3.93", "3.75"],
        ["คะแนนดาวเฉลี่ย", "3.09", "3.22", "2.99"],
    ], M, y + Inches(1.10), lw, [1.8, 1.5, 1.4, 1.4], row_h=Inches(0.30),
        header_size=9.8, body_size=9.8,
        aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])

    tf = textbox(slide, M, y + Inches(2.74), lw, Inches(1.6))
    bullets(tf, [
        "เตียงลดเพียง **5.1%** แต่ผู้พักอาศัยลด **17.4%** แล้วกลับมา **14.6%** "
        "— ตัวตั้งเป็นตัวที่เคลื่อน ไม่ใช่ตัวหาร การฟื้นจึงเป็นอุปสงค์ที่กลับมาจริง",
        "ชั่วโมงพยาบาลพุ่งเป็น 3.93 ตอนก้นเหว เพราะผู้พักอาศัยออกเร็วกว่าพนักงาน แล้วดิ่งเหลือ "
        "**3.61 ในปี 2566** ซึ่งคือวิกฤตขาดแคลนพยาบาลจริง — ปัจจุบันกลับมา 3.75 เท่าปี 2562 พอดี",
        "ข้อบกพร่องเพิ่มจาก 7.5 เป็น **9.2 ต่อแห่ง (+23.7%)**",
    ], size=10.4, gap=6, first=True)

    callout(slide, M, BODY_BOTTOM - Inches(1.34), lw, Inches(1.34),
            "ข้อนี้แก้คำตอบเดิม และนี่คือบทเรียนของโครงงาน",
            "ตอนที่คลังมีแค่ 4 งวด เราสรุปว่า “อัตราการเข้าพักฟื้นเพราะบ้านว่างปิดตัวไป” "
            "— **สรุปนั้นผิด** เพราะการเทียบสองปลายบังคับให้ต้องเดากลไกตรงกลาง "
            "ปลายทั้งสองบังเอิญใกล้กัน (79.8% กับ 80.2%) จนดูเหมือนแทบไม่มีอะไรเกิดขึ้น "
            "ทั้งที่ระหว่างนั้นมีการดิ่ง 13 จุดและการฟื้นห้าปี "
            "**ข้อมูลน้อยจุดไม่ได้ทำให้คำตอบกว้างขึ้น มันทำให้คำตอบผิดได้**",
            accent=WARN, fill=RGBColor(0xFB, 0xEE, 0xEB), title_size=10.5, body_size=9.8)

    place_image(slide, DASH / "bq5_trend.png", M + lw + Inches(0.3), y,
                CONTENT_W - lw - Inches(0.3), BODY_H)


def slide_bq6_bq7(prs):
    slide = new_slide(prs)
    y = head(slide, "03 · BQ6 + BQ7 · ความเสี่ยง",
             "แห่งใดกำลังเสื่อม และค่าปรับกระจุกตัวที่ใด")
    w = int((CONTENT_W - Inches(0.34)) / 2)
    place_image(slide, DASH / "bq6_risk_backtest.png", M, y, w, Inches(2.86))
    place_image(slide, DASH / "bq7_fines_timeline.png", M + w + Inches(0.34), y, w,
                Inches(2.86))

    ty = y + Inches(2.98)
    tf = textbox(slide, M, ty, w, Inches(2.4))
    para(tf, "BQ6 — คะแนนเสี่ยง 5 องค์ประกอบ ผ่านการทดสอบย้อนหลัง 4 ปี",
         size=12.5, color=ACCENT, font=HEAD_FONT, bold=True, first=True, space_after=6)
    bullets(tf, [
        "ให้คะแนนจากงวด **ม.ค. 2562** แล้วไปดูผลจริงปี 2566–2569 คะแนนไม่เคยเห็นผลลัพธ์เหล่านั้น",
        "คะแนน 0 → ถูกปรับ 37.1% · $204/เตียง · **3.45 ดาว**   ·   "
        "คะแนน 6 → ถูกปรับ 62.5% · $530/เตียง · **2.21 ดาว** — ขึ้นทุกขั้นในทุกคอลัมน์ผลลัพธ์",
        "กลุ่มเสี่ยงสูง (คะแนน ≥5) ถูกปรับ **1.30 เท่า** และจ่ายค่าปรับต่อเตียง **1.56 เท่า** ของฐาน",
        "สิ่งที่สำคัญกว่าตัวคูณ: บ้านที่เสี่ยงสูงในปี 2562 ยังได้เพียง 2.21 ดาวหลังผ่านไปเจ็ดปี "
        "— **ความเสี่ยงคงทน ไม่ใช่เรื่องชั่วคราว**",
    ], size=10.3, gap=6)

    tf = textbox(slide, M + w + Inches(0.34), ty, w, Inches(2.4))
    para(tf, "BQ7 — การบังคับใช้กลับทิศสองครั้ง ไม่ใช่แนวโน้มเดียว",
         size=12.5, color=ACCENT, font=HEAD_FONT, bold=True, first=True, space_after=6)
    bullets(tf, [
        "ปี 2562 ปรับ 3,126 ครั้ง เฉลี่ย $36,773 → ปี 2564 **17,370 ครั้ง เฉลี่ย $12,189** "
        "(ถี่ขึ้น 5.5 เท่า เบาลง 3 เท่า — คลื่นคุมการติดเชื้อ) → ปี 2568 กลับมา 3,894 ครั้ง เฉลี่ย $40,376",
        "ปี 2566 จึงเป็น**ขาขึ้นของการกลับสู่ภาวะปกติ** ไม่ใช่จุดเริ่มของแนวโน้มใหม่",
        "รายรัฐต่างกัน **เกือบ 15 เท่า** ใต้กฎกลางชุดเดียวกัน — VT $1,185/เตียง · IL $935 "
        "(และใหญ่ที่สุด $79.5M) · NH $81 · ระดับประเทศ $317",
        "**5 รัฐ (IL, TX, CA, OH, FL) กินค่าปรับ 43.0%** บนเตียงเพียง 31.8%",
    ], size=10.3, gap=6)


def slide_recommendations(prs):
    slide = new_slide(prs)
    y = head(slide, "03 · INSIGHTS & RECOMMENDATIONS",
             "ข้อเสนอแนะทางธุรกิจ 8 ข้อ (โจทย์กำหนดขั้นต่ำ 5 ข้อ)",
             sub="หนึ่งข้อต่อหนึ่งคำถาม ทุกข้อผูกกับตัวเลขบนแดชบอร์ด และถูกคำนวณใหม่เมื่อผู้ใช้ขยับตัวกรอง")
    items = [
        ("BQ1", "ชอร์ตลิสต์ New Mexico, Virginia, South Carolina และเลี่ยง Arkansas, Indiana, Iowa",
         "อุปทานต่ำ + เข้าพักสูง + คู่แข่งอ่อน ครบพร้อมกันทั้งสามรัฐแรก ส่วนสามรัฐหลังคือตลาดที่เตียงล้นแล้ว"),
        ("BQ2", "ตั้งราคาความเสี่ยงของ for-profit ในดีลซื้อกิจการ ไม่ถือเป็นคุณลักษณะที่เป็นกลาง",
         "ค่าปรับต่อเตียงสูงกว่า non-profit 1.56 เท่า และช่องว่างดาว 0.77 มาพร้อมชั่วโมงพยาบาลที่น้อยกว่า 16%"),
        ("BQ3", "งบกำลังคนก้อนถัดไปให้ 5,134 แห่งที่ยังต่ำกว่า 3.48 ชม. ไม่ใช่บ้านที่ผ่านเกณฑ์แล้ว",
         "ใต้เส้นค่าปรับสูงกว่า 1.51 เท่า แต่เหนือ 3.75 เส้นค่าปรับแบน เงินก้อนเดียวกันซื้อการลดความเสี่ยงไม่ได้อีก"),
        ("BQ4", "HR โจมตีการลาออกทุกระดับ · CFO ป้องเส้น 40% ซึ่งเป็นจุดที่เงินอยู่",
         "คุณภาพไม่มีระดับปลอดภัย แต่อัตราการเข้าพักหักที่ 40% พอดี และการรักษาคนเป็นคานงัดอิสระจากการจ้างเพิ่ม"),
        ("BQ5", "อย่าอ่านอัตราการเข้าพักที่ฟื้นว่าเป็นสัญญาณให้สร้างเตียงใหม่ทั่วประเทศ",
         "อุปสงค์กลับมาจริง แต่ผู้พักอาศัยยังต่ำกว่ายอดเดิม 5.3% และคุณภาพยังไม่ฟื้น ตัดสินใจที่ระดับรัฐ"),
        ("BQ6", "การซื้อกิจการที่ติดธงต้องแนบงบพลิกฟื้นมาด้วย ไม่ใช่หวังว่าจะดีขึ้นเอง",
         "บ้านที่เสี่ยงสูงในปี 2562 ยังได้เพียง 2.21 ดาวในปี 2569 — ความเสี่ยงคงทนข้ามเจ็ดปี"),
        ("BQ7", "ถ่วงงบกำกับกฎระเบียบตามรัฐ ไม่ใช่ตามจำนวนสถานพยาบาล และวางแผนรับเหตุการณ์ใหญ่จำนวนน้อย",
         "รายรัฐต่างกันเกือบ 15 เท่าใต้กฎกลางชุดเดียวกัน และ 5 รัฐกินค่าปรับ 43.0% บนเตียงเพียง 31.8%"),
        ("BQ8", "อย่าจ่ายพรีเมียมด้านคุณภาพให้ขนาดเครือ — จ่ายตามข้อได้เปรียบด้านอัตราการเข้าพัก",
         "เครือ 50+ แห่งได้เข้าพักสูงกว่า 1.6 จุด แต่ชั่วโมงพยาบาลน้อยกว่า 13.8% และดาวต่ำกว่า 0.32"),
    ]
    cw = int((CONTENT_W - Inches(0.28)) / 2)
    for i, (bq, action, why) in enumerate(items):
        col, row = i % 2, i // 2
        x = M + col * (cw + Inches(0.28))
        yy = y + row * Inches(1.14)
        rect(slide, x, yy, cw, Inches(1.02), fill=WHITE, line=GRID)
        rect(slide, x, yy, Inches(0.05), Inches(1.02), fill=GREEN)
        tag = textbox(slide, x + Inches(0.22), yy + Inches(0.14), Inches(0.6), Inches(0.3))
        _rich(tag.paragraphs[0], bq, size=11, color=GREEN, font=HEAD_FONT, bold=True)
        tf = textbox(slide, x + Inches(0.86), yy + Inches(0.13), cw - Inches(1.1), Inches(0.8))
        para(tf, action, size=11.2, color=INK, bold=True, first=True, space_after=3,
             line_spacing=1.1)
        para(tf, why, size=9.4, color=MUTED, space_after=0, line_spacing=1.1)


def slide_ai(prs):
    slide = new_slide(prs)
    y = head(slide, "03 · GENERATIVE AI", "ใช้ AI อย่างไร และตรวจว่าผลลัพธ์ถูกต้องอย่างไร",
             sub="บันทึกไว้ใน 05_AI_Usage_Log/AI_Usage_Log.md — ปัจจุบัน 18 รายการ จากขั้นต่ำที่โจทย์กำหนด 5 รายการ")
    cw = int((CONTENT_W - Inches(0.3)) / 2)
    callout(slide, M, y, cw, Inches(1.68), "อะไรนับว่า “ตรวจแล้ว”",
            "รันคำสั่งแล้วเทียบกับค่าที่รู้คำตอบอยู่ก่อน · เปิดแฟ้มข้อมูลจริงแล้วนับแถวเอง · "
            "เทียบกับพจนานุกรมข้อมูลของ CMS · หรือให้สมาชิกอีกคนอ่านงานแล้วอธิบายกลับมาได้",
            body_size=11)
    callout(slide, M + cw + Inches(0.3), y, cw, Inches(1.68), "อะไรไม่นับ",
            "“อ่านแล้วดูสมเหตุสมผล” · “AI บอกว่ามั่นใจ” · “ถามซ้ำแล้วได้คำตอบเดิม” · "
            "“รันแล้วไม่ error” — ทั้งสี่ข้อนี้ไม่ใช่การตรวจสอบ",
            accent=WARN, fill=RGBColor(0xFB, 0xEE, 0xEB), body_size=11)

    ty = y + Inches(1.94)
    tf = textbox(slide, M, ty, CONTENT_W, Inches(0.3))
    para(tf, "สามตัวอย่างที่ผลลัพธ์จาก AI **ไม่ได้** ถูกรับมาใช้ทั้งดุ้น",
         size=13, color=INK, font=HEAD_FONT, bold=True, first=True, space_after=0)
    cases = [
        ("ปฏิเสธทันที", WARN,
         "AI เสนอให้เพิ่ม payment_denial_start_date เข้าไปในคีย์กำจัดของซ้ำของ Fact 2 "
         "มันแยกแยะได้ดีขึ้นจริง แต่คุณลักษณะที่ CMS แก้ย้อนหลังได้ต้องไม่อยู่ในคีย์ "
         "— ถ้าต้นทางแก้วันที่ ค่าปรับรายการเดิมจะกลายเป็นแถวใหม่ทันที คีย์จึงถูกเก็บให้แคบไว้"),
        ("ปฏิเสธเชิงหลักการ", ORANGE,
         "สองข้อเสนอถูกปฏิเสธด้วยเหตุผลเดียวกัน คือการซ่อนข้อบกพร่องแย่กว่าการแสดงมัน "
         "— การตัดแถว is_suspect ออกจาก view โดยปริยายขัดกับกฎ “ติดธง ไม่ลบ” จึงรายงาน "
         "suspect_rows ควบคู่ไปแทน · และการซ่อนแท็บ BQ1 ระหว่างที่ M10 ยังไม่ครบ "
         "จะทำให้ผู้อ่านไม่มีวันรู้ว่ามีช่องว่าง แท็บจึงอยู่ต่อพร้อมกล่องแดงบอกว่าอะไรขาด"),
        ("รับไว้ แต่ตรวจเองก่อน", GREEN,
         "AI รายงานว่า pipeline ทำแถวค่าปรับหายไปหนึ่งในสาม เราไม่เชื่อทันที "
         "— ไล่รายชื่อแฟ้มที่แตกออกมาของทั้งสองงวด นับแถว CSV ใหม่ด้วย pandas (16,180 และ 16,166) "
         "ทำยอดเดิมขึ้นมาเองด้วยมือ (8,789 + 16,180 + 16,166 = 41,135) แล้วไล่หา commit ต้นเหตุด้วย git show "
         "จึงค่อยบันทึกว่าเป็นข้อบกพร่องจริง · ส่วนข้อเสนอให้ revert นั้นถูกปฏิเสธ เพราะเหตุผลเดิมถูก และวิธีแก้ที่ถูกต้องแคบกว่านั้น"),
    ]
    cw3 = int((CONTENT_W - Inches(0.36)) / 3)
    for i, (title, col, body) in enumerate(cases):
        x = M + i * (cw3 + Inches(0.18))
        yy = ty + Inches(0.42)
        rect(slide, x, yy, cw3, Inches(2.42), fill=WHITE, line=GRID)
        rect(slide, x, yy, cw3, Inches(0.05), fill=col)
        tf = textbox(slide, x + Inches(0.24), yy + Inches(0.22), cw3 - Inches(0.48),
                     Inches(2.0))
        para(tf, title, size=12.5, color=col, font=HEAD_FONT, bold=True, first=True,
             space_after=5)
        para(tf, body, size=10, color=INK2, space_after=0, line_spacing=1.15)


def slide_limitations(prs):
    slide = new_slide(prs)
    y = head(slide, "ข้อจำกัด", "สิ่งที่ยังทำไม่ได้ — ประกาศเอง ไม่ปล่อยให้ผู้อ่านไปเจอเอง")
    items = [
        ("คะแนนดาวเป็นไม้บรรทัดที่ขยับเอง",
         "CMS ปรับเกณฑ์การให้ดาวระหว่างทาง การตกจาก 3.35 เหลือ 2.99 จึงปนการเปลี่ยนเกณฑ์กับ"
         "การเปลี่ยนคุณภาพจริง ข้อสรุปเรื่องคุณภาพใน BQ5 จึงยืนบน**จำนวนข้อบกพร่อง** (+23.7%) เป็นหลัก"),
        ("ข้อมูลเครือมีเฉพาะปี 2569",
         "คอลัมน์เครือเพิ่งเพิ่มเข้ามาปี 2568 BQ8 จึงยังแยกไม่ได้ว่า “เครือทำให้บ้านแย่ลง” "
         "หรือ “เครือไปซื้อบ้านที่แย่อยู่แล้ว” · ค่าปรับของยุคก่อน 2566 ชี้ไปที่แถว Unknown ไม่ใช่ “ไม่สังกัดเครือ”"),
        ("อัตราการลาออกมีเฉพาะยุค 2569",
         "BQ4 จึงเป็น cross-sectional ไม่ใช่ longitudinal — บอกได้ว่าอัตราลาออกสูงไปด้วยกันกับ"
         "คุณภาพต่ำ แต่บอกไม่ได้ว่าอันไหนมาก่อน"),
        ("ประชากรมีสองเส้นทางที่ยังไม่รวมกัน",
         "Ref_State_Population ใช้งานได้และเป็นที่มาของตัวเลข BQ1 ทั้งหมด แต่คอลัมน์เดิม "
         "Dim_Geography.pop_65plus ยังว่าง view v_market_saturation และแดชบอร์ดชุด Plotly "
         "จึงยังคืนค่า null สำหรับ M10 — งานที่เหลือคือย้าย view มาอ่าน Ref_State_Population"),
        ("แหล่ง VBP ยังไม่ถูกทำแบบจำลอง",
         "S3 ถูก extract แล้วแต่ยังไม่มีโมดูลใดอ่าน ตัวคูณจูงใจจึงยังไม่ถึงคลัง "
         "· Fact ตัวที่สามที่ grain เป็น (สถานพยาบาล × ปีงบประมาณ) ออกแบบไว้แล้วแต่ยังไม่ได้สร้าง"),
        ("รายไตรมาส ไม่ใช่รายเดือน",
         "CMS มีครบ 88 งวดรายเดือน โครงงานใช้ 32 งวดรายไตรมาสเพื่อคุมขนาดข้อมูล "
         "ถ้าต้องการละเอียดกว่านี้ รัน fetch_snapshots.py --all (~3 GB) แล้ว build ใหม่ได้ทันที"),
    ]
    cw = int((CONTENT_W - Inches(0.28)) / 2)
    for i, (t, b) in enumerate(items):
        col, row = i % 2, i // 2
        x = M + col * (cw + Inches(0.28))
        yy = y + row * Inches(1.62)
        rect(slide, x, yy, cw, Inches(1.46), fill=WHITE, line=GRID)
        rect(slide, x, yy, Inches(0.05), Inches(1.46), fill=MUTED)
        tf = textbox(slide, x + Inches(0.26), yy + Inches(0.18), cw - Inches(0.5),
                     Inches(1.12))
        para(tf, "%d. %s" % (i + 1, t), size=12, color=INK, font=HEAD_FONT, bold=True,
             first=True, space_after=4)
        para(tf, b, size=10, color=INK2, space_after=0, line_spacing=1.14)


def slide_closing(prs):
    slide = new_slide(prs, bg=INK)
    rect(slide, 0, 0, Inches(0.09), SLIDE_H, fill=ACCENT)
    tf = textbox(slide, Inches(1.15), Inches(0.85), Inches(11), Inches(0.9))
    para(tf, "สรุป", size=13, color=ACCENT, font=HEAD_FONT, bold=True, first=True,
         space_after=6, line_spacing=1.0)
    para(tf, "จากแฟ้ม ZIP 88 แฟ้มที่ไม่มีใครต่อกัน สู่คำตอบที่ตัดสินใจได้ 8 ข้อ",
         size=28, color=WHITE, font=HEAD_FONT, bold=True, space_after=0, line_spacing=1.05)

    stats = [("5", "แหล่งข้อมูล", "3 แหล่งดึงผ่าน HTTP"),
             ("8", "กฎคุณภาพข้อมูล", "1,555 แถวใน etl_run_log"),
             ("2 + 6", "Fact + Dimension", "483,183 + 79,803 แถว"),
             ("115", "การตรวจที่ผ่าน", "24 + 24 + 56 + 11")]
    w = int((Inches(11.0) - Inches(0.6)) / 4)
    for i, (v, l, n) in enumerate(stats):
        x = Inches(1.15) + i * (w + Inches(0.2))
        tf = textbox(slide, x, Inches(2.42), w, Inches(1.1))
        para(tf, v, size=30, color=ACCENT, font=HEAD_FONT, bold=True, first=True,
             space_after=2, line_spacing=1.0)
        para(tf, l, size=12, color=WHITE, space_after=1)
        para(tf, n, size=10, color=MUTED, space_after=0)

    hline(slide, Inches(1.15), Inches(3.86), Inches(11.0), color=RGBColor(0x3A, 0x39, 0x36))

    tf = textbox(slide, Inches(1.15), Inches(4.08), Inches(6.0), Inches(2.2))
    para(tf, "คะแนนพิเศษที่เข้าเกณฑ์", size=12.5, color=ACCENT, font=HEAD_FONT, bold=True,
         first=True, space_after=7)
    bullets(tf, [
        "ดึงข้อมูลผ่าน **API** สองแหล่ง และประกอบ **nested JSON** ของสำมะโนประชากร",
        "แดชบอร์ดด้วย **Streamlit** พร้อม drill-down รายรัฐและรายสถานพยาบาล",
        "**Incremental load** ในชั้น ETL และ pipeline ที่รันซ้ำได้ทั้งสาย",
        "**ตรวจคุณภาพข้อมูลอัตโนมัติ** ที่อ่าน etl_run_log สด ๆ บนหน้าจอ",
        "เตรียม **analytical dataset** สำหรับต่อยอด (คะแนนเสี่ยงที่ผ่าน backtest 4 ปี)",
    ], size=10.8, gap=6, color=RGBColor(0xC8, 0xC6, 0xC0), bullet_color=ACCENT)

    tf = textbox(slide, Inches(7.55), Inches(4.08), Inches(4.6), Inches(2.2))
    para(tf, "ลิงก์", size=12.5, color=ACCENT, font=HEAD_FONT, bold=True, first=True,
         space_after=7)
    for label in ("โฟลเดอร์ Google Drive (Anyone with the link can view)",
                  "____________________________________________",
                  "แดชบอร์ด",
                  "____________________________________________",
                  "รายงานฉบับเต็ม  06_Report/eldercare_report.pdf"):
        para(tf, label, size=10.8, color=RGBColor(0xC8, 0xC6, 0xC0), space_after=5)


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    slide_title(prs)
    slide_agenda(prs)

    slide_section(prs, "ส่วนที่ 01  ·  2 นาที", "Business Problem\n& Requirements",
                  "ปัญหาทางธุรกิจ · Stakeholders · Business Questions · Measures", ACCENT)
    slide_problem(prs)
    slide_stakeholders(prs)
    slide_questions(prs)
    slide_measures(prs)
    slide_ratio_trap(prs)

    slide_section(prs, "ส่วนที่ 02  ·  5 นาที", "Data Warehouse\nDesign & ETL",
                  "แหล่งข้อมูล · คุณภาพข้อมูล · Star Schema · Grain · ETL · การตรวจสอบ", ORANGE)
    slide_sources(prs)
    slide_quality(prs)
    slide_star(prs)
    slide_grain(prs)
    slide_built(prs)
    slide_etl(prs)
    slide_verification(prs)

    slide_section(prs, "ส่วนที่ 03  ·  5 นาที", "Dashboard\n& Business Insights",
                  "สาธิตแดชบอร์ด · คำตอบ BQ1–BQ8 · ข้อเสนอแนะ · การใช้ Generative AI", GREEN)
    slide_dashboard(prs)
    slide_filters(prs)
    slide_bq1(prs)
    slide_bq2_bq8(prs)
    slide_bq3(prs)
    slide_bq4(prs)
    slide_bq5(prs)
    slide_bq6_bq7(prs)
    slide_recommendations(prs)
    slide_ai(prs)
    slide_limitations(prs)
    slide_closing(prs)

    # ใส่เลขหน้าให้ทุกสไลด์ยกเว้นหน้าปกและหน้าปิด ซึ่งเป็นพื้นเข้มและไม่ต้องการ footer
    for i, slide in enumerate(prs.slides):
        if i in (0, len(prs.slides._sldIdLst) - 1):
            continue
        _page[0] = i
        footer(slide)

    prs.save(OUT)
    print("เขียนแล้ว: %s  (%d สไลด์)" % (OUT, len(prs.slides._sldIdLst)))


if __name__ == "__main__":
    main()
